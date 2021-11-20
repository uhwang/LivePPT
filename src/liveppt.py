'''
    Author: Uisang Hwang
    Email : uhwangtx@gmail.com

    LivePPT Ver 0.1
    
    08/07/20  Custom Image Resolution @ PPP to Image
              Custom Shadow effect
    08/08/20  Separation of outline and shadow effect 
              Each presentation opens a same file.
              Add Text align
    08/13/20  Merge ppt files
    08/15/20  Add dash, transprancy in Fx outline
    08/21/20  Add Gradient fill in a text frame
              Add margin (L,T,R,B) of a text frame 
    08/23/20  Rewrite the code
    09/30/20  Remove tabs
    11/17/20  Add tooptip & Path copy button(src to save) 
    12/14/20  Add custom worship type dialogue
    01/07/21  Add Responsive reading
    01/09/21  Add Responsive reading in range(1-137)
              Bug fix
              Add Responsive reading #136
    01/11/21  Add Responsive reading number range
    
    02/14/21  1) new hymnal database from Bathlehtm 4.3.1
              2) rewrite hymnal parse function for corus at hymal.py
              3) switch from cx_Freeze to Py2Exe
                 ImportError: (cannot import name _elementpath) xml.etree.pyd'
                 Add 'lxml.etree', 'lxml._elementpath' to 'include'
              4) Add default.pptx to Document for creating pptx
    03/01/21  Rewrite hymal.py for chorus and create_hymal_ppt for numbering
    03/02/21  File deleting exception
    05/04/21  Extract Txt from PPTX
    11/20/21  Separate Slide Tab as a dialogue box

    Convert praise ppt to subtitle ppt for live streaming

    Requirements
    ----------------------------------------------
    * Read multiple praise ppt
    * Get slide size (4:3 or 16:9)
    * Get text box size: sx, sy, wid, hgt
    * Get text properties: typeface, size(point), color (custom, rgb)
    * Insert blank slide after a praise ppt
    * Can change the order of source ppt files
    * Can delete a source ppt file
    * Remove all ppt files from list box

    Dependencies
    ----------------------------------------------
    * PyQt4
    * Python pptx
    * win32com

    Outline text VBA code
    ----------------------------------------------
    Sub outline()
    Dim sld As Slide
    Dim shp As Shape

    For Each sld In ActivePresentation.Slides
        sld.Select
        ActiveWindow.ViewType = ppViewSlide
        ActiveWindow.Activate

        sld.Shapes(1).Select
        With ActiveWindow.Selection.TextRange2.Font
            .Line.Visible = msoCTrue
            .Line.Pattern = msoPattern10Percent
            .Line.ForeColor.RGB = RGB(255, 0, 0)
            .Line.Weight = 2
            .Line.Style = msoLineSingle
            .Line.DashStyle = msoLineDash
            .Line.Transparency = 1
        End With
    Next sld
    End Sub

    Shadow text VBA code
    ----------------------------------------------
    Sub shadow()
    Dim sld As Slide
    Dim shp As Shape
    For Each sld In ActivePresentation.Slides
        For Each shp In sld.Shapes
            With shp.TextFrame2.TextRange.Font
            .shadow.Visible = msoCTrue
            .shadow.Style = msoShadowStyleOuterShadow
            .shadow.OffsetX = 2
            .shadow.OffsetY = 2
            ' !!!! Do not use size option !!!!
            '.shadow.Size = 1
            .shadow.Blur = 2
            .shadow.Transparency = 0.7
            End With
        Next shp
    Next sld
    End Sub
'''

import re
import os
import datetime
import pptx
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import sys
from PyQt4 import QtCore, QtGui
import msoLine
import msoDash
import msoShadow
import hymal
import msgcom
import livepptconst as const
import livepptcolor as color
import livepptfunc as func
from livepptcls import *
import liveppt_slidedata_dlg as slide_data_dlg

try:
    import win32com.client
except:
    pass

import icon_file_add
import icon_folder_open
import icon_arrow_down
import icon_arrow_up     
import icon_delete        
import icon_table_sort_asc  
import icon_table_sort_desc 
import icon_trash
import icon_convert
import icon_color_picker
import icon_font_picker
import icon_ppt_pptx
import icon_ppt_image
import icon_merge
import icon_shadow
import icon_outline
import icon_liveppt
import icon_copy_src_path
import icon_keyboard
import icon_respread
import icon_respread_sub
import icon_clear
import icon_restore
import icon_ppt_txt
import icon_slide_setting
import icon_fit

class QBibleChecker(QtGui.QDialog):
    def __init__(self, lppt):
        super(QBibleChecker, self).__init__()
        self.initUI()

    def check_db_tab_UI(self):
        layout = QtGui.QFormLayout()
    
        kbib_group = QtGui.QGroupBox('Korean Bible')
        kbib_layout = QtGui.QGridLayout()
        kcheck_list = bwxref.get_kbible_check_list()
        self.kbible_checker = {}
        ncol = 3
        row = 0
    
        for i, kbib in enumerate(bwxref.get_kbible_list()):
            checker = QtGui.QCheckBox(kbib, self)
            checker.setChecked(kcheck_list[kbib])
            self.kbible_checker[kbib] = checker
            col = i % ncol
            row = row+1 if i is not 0 and i%ncol is 0 else row
            kbib_layout.addWidget(self.kbible_checker[kbib], row, col)
            
        kbib_group.setLayout(kbib_layout)
        layout.addRow(kbib_group)		
        
        ebib_group = QtGui.QGroupBox('English Bible')
        ebib_layout = QtGui.QGridLayout()
        echeck_list = bwxref.get_ebible_check_list()
        self.ebible_checker = {}
        ncol = 3
        row = 0
    
        for i, ebib in enumerate(bwxref.get_ebible_list()):
            checker = QtGui.QCheckBox(ebib, self)
            checker.setChecked(echeck_list[ebib])
            self.ebible_checker[ebib] = checker
            col = i % ncol
            row = row+1 if i is not 0 and i%ncol is 0 else row
            ebib_layout.addWidget(self.ebible_checker[ebib], row, col)
            
        ebib_group.setLayout(ebib_layout)
        layout.addRow(ebib_group)		
        
        
        hgbib_group = QtGui.QGroupBox('Hebrew/Greek Bible')
        hgbib_layout = QtGui.QGridLayout()
        hgcheck_list = bwxref.get_hgbible_check_list()
        self.hgbible_checker = {}
        ncol = 3
        row = 0
        for i, hgbib in enumerate(bwxref.get_hgbible_list()):
            checker = QtGui.QCheckBox(hgbib, self)
            checker.setChecked(hgcheck_list[hgbib])
            self.hgbible_checker[hgbib] = checker
            #checker.setEnabled(True)
            col = i % ncol
            row = row+1 if i is not 0 and i%ncol is 0 else row
            hgbib_layout.addWidget(self.hgbible_checker[hgbib], row, col)
        hgbib_group.setLayout(hgbib_layout)  
        layout.addRow(hgbib_group)	
        
        self.check_db_tab.setLayout(layout)
    
    
    
class QUserWorshipType(QtGui.QDialog):
    def __init__(self, lppt):
        super(QUserWorshipType, self).__init__()
        self.lppt = lppt
        self.initUI()
        
    def initUI(self):
        layout = QtGui.QFormLayout()
        # Create an array of radio buttons
        moods = [QtGui.QRadioButton("Current"), QtGui.QRadioButton("User")]

        # Radio buttons usually are in a vertical layout   
        source_layout = QtGui.QHBoxLayout()

        # Create a button group for radio buttons
        self.mood_button_group = QtGui.QButtonGroup()

        for i in range(len(moods)):
            # Add each radio button to the button layout
            source_layout.addWidget(moods[i])
            # Add each radio button to the button group & give it an ID of i
            self.mood_button_group.addButton(moods[i], i)
            # Connect each radio button to a method to run when it's clicked
            self.connect(moods[i], QtCore.SIGNAL("clicked()"), self.radio_button_clicked)

        # Set a radio button to be checked by default
        moods[0].setChecked(True)   
        
        source_type_layout = QtGui.QVBoxLayout()
        self.ppt_list = QtGui.QComboBox()
        self.ppt_list.addItems(self.lppt.get_ppt_list())
        self.user_input = QtGui.QLineEdit()
        source_type_layout.addWidget(self.ppt_list)
        source_type_layout.addWidget(self.user_input)
        
        button_layout = QtGui.QHBoxLayout()
        self.ok = QtGui.QPushButton('OK')
        self.ok.clicked.connect(self.accept)
        button_layout.addWidget(self.ok)

        self.no = QtGui.QPushButton('Cancel')
        self.no.clicked.connect(self.reject)
        button_layout.addWidget(self.no)

        layout.addRow(source_layout)
        layout.addRow(source_type_layout)
        layout.addRow(button_layout)
        self.setLayout(layout)
        self.setWindowTitle('Custom Worship Type')
        self.radio_button_clicked()
        
    def radio_button_clicked(self):
        id = self.mood_button_group.checkedId()
        if id == 0: # get output name from ppt files
            self.ppt_list.setEnabled(True)
            self.user_input.setEnabled(False)
        elif id == 1: # get output name from user input
            self.ppt_list.setEnabled(False)
            self.user_input.setEnabled(True)
    
    def get_source(self):
        return self.mood_button_group.checkedId()
        
class QShadowInfo(QtGui.QDialog):
    def __init__(self, shd):
        super(QShadowInfo, self).__init__()
        self.initUI(shd)

    def initUI(self, shd):
        layout = QtGui.QFormLayout()
        button_layout = QtGui.QHBoxLayout()
        item_layout = QtGui.QVBoxLayout()

        self.style = QtGui.QComboBox(self)
        shadow_style = ["Inner", "Outer", "Mixed"]
        self.style.addItems(shadow_style)
        self.style.setCurrentIndex(shd.Style-1)
        item_layout.addWidget(self.style)
        item_layout.addWidget(QtGui.QLabel("Offset(X)"))
        self.offset_x = QtGui.QLineEdit("%d"%shd.OffsetX)
        item_layout.addWidget(self.offset_x)
        item_layout.addWidget(QtGui.QLabel("Offset(Y)"))
        self.offset_y = QtGui.QLineEdit("%d"%shd.OffsetY)
        item_layout.addWidget(self.offset_y)
        item_layout.addWidget(QtGui.QLabel("Blur"))
        self.blur = QtGui.QLineEdit("%d"%shd.Blur)
        item_layout.addWidget(self.blur)
        item_layout.addWidget(QtGui.QLabel("Transparency"))
        self.trans = QtGui.QLineEdit("%f"%shd.Transparency)
        item_layout.addWidget(self.trans)

        self.ok = QtGui.QPushButton('OK')
        self.ok.clicked.connect(self.accept)
        button_layout.addWidget(self.ok)

        self.no = QtGui.QPushButton('Cancel')
        self.no.clicked.connect(self.reject)
        button_layout.addWidget(self.no)

        layout.addRow(item_layout)
        layout.addRow(button_layout)
        self.setLayout(layout)

    def get_shadow_info(self):
        st = int(self.style.currentIndex())
        ox = int(self.offset_x.text())
        oy = int(self.offset_y.text())
        bl = int(self.blur.text())
        tr = float(self.trans.text())
        return st, ox, oy, bl, tr

class QImageResolution(QtGui.QDialog):
    def __init__(self):
        super(QImageResolution, self).__init__()
        self.initUI()

    def initUI(self):
        layout = QtGui.QFormLayout()
        button_layout = QtGui.QHBoxLayout()
        item_layout = QtGui.QHBoxLayout()

        self.res = QtGui.QComboBox(self)
        self.res_list = ["HD:1280x720", "Full HD:1920x1080",\
                         "Quad HD:2560x1440", "Ultra HD:3840x2160"]
        self.res.addItems(self.res_list)
        item_layout.addWidget(self.res)
        self.res.setCurrentIndex(0)

        self.ok = QtGui.QPushButton('OK')
        self.ok.clicked.connect(self.accept)
        button_layout.addWidget(self.ok)

        self.no = QtGui.QPushButton('Cancel')
        self.no.clicked.connect(self.reject)
        button_layout.addWidget(self.no)

        layout.addRow(item_layout)
        layout.addRow(button_layout)
        self.setLayout(layout)

    def get_resolution(self):
        r = self.res.currentText().split(':')[1].split('x')
        return int(r[0]), int(r[1])

class QChooseFxSource(QtGui.QDialog):
    def __init__(self):
        super(QChooseFxSource, self).__init__()
        self.initUI()

    def initUI(self):
        layout = QtGui.QFormLayout()
        # Create an array of radio buttons
        moods = [QtGui.QRadioButton("Current"), QtGui.QRadioButton("Fx")]

        # Set a radio button to be checked by default
        moods[1].setChecked(True)   

        # Radio buttons usually are in a vertical layout   
        source_layout = QtGui.QHBoxLayout()

        # Create a button group for radio buttons
        self.mood_button_group = QtGui.QButtonGroup()

        for i in range(len(moods)):
            # Add each radio button to the button layout
            source_layout.addWidget(moods[i])
            # Add each radio button to the button group & give it an ID of i
            self.mood_button_group.addButton(moods[i], i)
            # Connect each radio button to a method to run when it's clicked
            #self.connect(moods[i], QtCore.SIGNAL("clicked()"), self.radio_button_clicked)

        self.radio_button_clicked()
        button_layout = QtGui.QHBoxLayout()
        self.ok = QtGui.QPushButton('OK')
        self.ok.clicked.connect(self.accept)
        button_layout.addWidget(self.ok)

        self.no = QtGui.QPushButton('Cancel')
        self.no.clicked.connect(self.reject)
        button_layout.addWidget(self.no)

        layout.addRow(source_layout)
        layout.addRow(button_layout)

        self.setLayout(layout)

    def radio_button_clicked(self):
        return

    def get_source(self):
        return self.mood_button_group.checkedId()

class QLivePPT(QtGui.QWidget):
    def __init__(self):
        super(QLivePPT, self).__init__()
        self.initUI()

    def initUI(self):
        self.set_common_var()
        self.form_layout  = QtGui.QFormLayout()
        tab_layout = QtGui.QVBoxLayout()
        self.tabs = QtGui.QTabWidget()
        policy = self.tabs.sizePolicy()
        policy.setVerticalStretch(1)
        self.tabs.setSizePolicy(policy)
        self.tabs.setEnabled(True)
        self.tabs.setTabPosition(QtGui.QTabWidget.West)
        self.tabs.setObjectName('Media List')

        self.ppt_tab = QtGui.QWidget()
        #self.slide_tab = QtGui.QWidget()
        self.hymal_tab = QtGui.QWidget()
        self.fx_tab = QtGui.QWidget()
        self.respread_tab = QtGui.QWidget()
        self.message_tab = QtGui.QWidget()
        self.txtppt_tab = QtGui.QWidget()
        self.bibppt_tab = QtGui.QWidget()
        
        self.tabs.addTab(self.ppt_tab     , const._ppttab_text)
        #self.tabs.addTab(self.slide_tab   , const._slidetab_text)
        self.tabs.addTab(self.fx_tab      , const._fxtab_text)
        self.tabs.addTab(self.hymal_tab   , const._hymaltab_text)
        self.tabs.addTab(self.respread_tab, const._respread_text)
        self.tabs.addTab(self.txtppt_tab  , const._txtppt_text)
        self.tabs.addTab(self.bibppt_tab  , const._bibppt_text)
        self.tabs.addTab(self.message_tab , const._messagetab_text)
        
        self.message_tab_UI()
        self.ppt_tab_UI()
        #self.slide_tab_UI()
        self.hymal_tab_UI()
        self.respreading_tab_UI()
        self.fx_tab_UI()
        self.txtppt_tab_UI()
        self.bibppt_tab_UI()
        
        tab_layout.addWidget(self.tabs)
        self.form_layout.addRow(tab_layout)
        self.setLayout(self.form_layout)

        self.setWindowTitle("LivePPT")
        self.setWindowIcon(QtGui.QIcon(QtGui.QPixmap(icon_liveppt.table)))
        self.show()

    def set_bibppt_slide_size(self):
        return
    
    def create_bible_ppt(self):
        return
        
    def change_bible_db_path(self):
        return
        
    def bibppt_tab_UI(self):
        layout = QtGui.QFormLayout()
            
        grid = QtGui.QGridLayout()
        grid.addWidget(QtGui.QLabel("Search"), 0, 0)
        self.bible_verse = QtGui.QLineEdit()
        self.bible_verse.setToolTip("창 1:1-2")
        grid.addWidget(self.bible_verse, 0, 1)
            
        grid.addWidget(QtGui.QLabel('Slide Size'), 1, 0)
        self.choose_bibppt_slide_size = QtGui.QComboBox()
        self.choose_bibppt_slide_size.addItems(const._slide_size_type)
        self.choose_bibppt_slide_size.setCurrentIndex(0)
        self.choose_bibppt_slide_size.currentIndexChanged.connect(self.set_bibppt_slide_size)
        grid.addWidget(self.choose_bibppt_slide_size, 1, 1)    
            
        grid.addWidget(QtGui.QLabel('Bible DB'), 2, 0)
        self.bible_db_file_path  = QtGui.QLineEdit(os.getcwd())
        self.bible_db_file_path_btn = QtGui.QPushButton('', self)
        self.bible_db_file_path_btn.clicked.connect(self.change_bible_db_path)
        self.bible_db_file_path_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.bible_db_file_path_btn.setIconSize(QtCore.QSize(16,16))
        #self.bible_db_file_path_btn.clicked.connect(self.get_slide_data)
        grid.addWidget(self.bible_db_file_path, 2, 1)
        grid.addWidget(self.bible_db_file_path_btn, 2, 2)
        
        layout.addRow(grid)
        
        self.bibppt_tab.setLayout(layout)

        #if res.exec_() is not 1:
        #    return
    '''
    responsive reading format
    text_box sx sy wid hgt
    title font_name font_size bold color align wrap
    newline font_size
    text font_name font_size bold color align wrap
    text font_name font_size bold color align wrap
    newline font_size
    text font_name font_size bold color align wrap
    text font_name font_size bold color align wrap
    '''
    
    def write_respread_format(self):
        with open(_responsive_reading_file, "wt") as w:
            for fm in _responsive_reading_default_format:
                w.write('%s\n'%fm)
                
    def respreading_tab_UI(self):
        layout = QtGui.QFormLayout()
        lay1 = QtGui.QHBoxLayout()   
        lay1.addWidget(QtGui.QLabel("R.Read Num"))
        self.respread_num = QtGui.QLineEdit("1")
        #self.respread_num.setValidator(QtGui.QIntValidator(1,137))
        lay1.addWidget(self.respread_num)
        
        self.respread_format_tbl = QtGui.QTableWidget()
        #font = QtGui.QFont("Fixedsys",9,True)
        #self.respread_format_tbl.setFont(font)
        self.respread_format_tbl.horizontalHeader().hide()
        self.respread_format_tbl.verticalHeader().hide()
        self.respread_format_tbl.setColumnCount(2)
        self.respread_format_tbl.setHorizontalHeaderItem(0, QtGui.QTableWidgetItem("Item"))
        self.respread_format_tbl.setHorizontalHeaderItem(1, QtGui.QTableWidgetItem("Data"))
        self.respread_format_tbl.setHorizontalHeaderItem(2, QtGui.QTableWidgetItem("Change"))

        header = self.respread_format_tbl.horizontalHeader()
        header.setResizeMode(0, QtGui.QHeaderView.ResizeToContents)
        header.setResizeMode(1, QtGui.QHeaderView.ResizeToContents)
        header.setResizeMode(2, QtGui.QHeaderView.ResizeToContents)        
        header.setResizeMode(1, QtGui.QHeaderView.Stretch)

        self.restore_responsive_reading_format()
        #info = ["TxtBox", "Title", "NewLine", "Text", "Text", "NewLine", "Text", "Text"]
        #_responsive_reading_format_key
        #self.respread_format_tbl.setRowCount(len(_responsive_reading_default_format))
        #for ii, jj in enumerate(_responsive_reading_default_format):
        #    item = QtGui.QTableWidgetItem(jj)
        #    item.setFlags(QtCore.Qt.ItemIsEnabled)
        #    fmt = _responsive_reading_default_format[ii].split('|')
        #    self.respread_format_tbl.setItem(ii, 0, QtGui.QTableWidgetItem(fmt[0]))
        #    self.respread_format_tbl.setItem(ii, 1, QtGui.QTableWidgetItem(fmt[1]))
        
        lay2 = QtGui.QGridLayout()
        lay2.addWidget(QtGui.QLabel('Slide Size'), 0, 0)
        self.choose_respread_slide_size = QtGui.QComboBox(self)
        self.choose_respread_slide_size.addItems(const._slide_size_type)
        self.choose_respread_slide_size.setCurrentIndex(0)
        self.choose_respread_slide_size.currentIndexChanged.connect(self.set_respread_slide_size)
        lay2.addWidget(self.choose_respread_slide_size, 0, 1)
        
        #lay3 = QtGui.QHBoxLayout()
        lay2.addWidget(QtGui.QLabel('Dest'), 1, 0)
        self.respread_dest_path  = QtGui.QLineEdit(os.getcwd())
        self.respread_dest_path_btn = QtGui.QPushButton('', self)
        self.respread_dest_path_btn.clicked.connect(self.get_respread_path)
        self.respread_dest_path_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.respread_dest_path_btn.setIconSize(QtCore.QSize(16,16))
        self.respread_dest_path_btn.setToolTip('Respread folder')
        lay2.addWidget(self.respread_dest_path, 1, 1)
        lay2.addWidget(self.respread_dest_path_btn, 1, 2)

        lay2.addWidget(QtGui.QLabel('File Check'), 2, 0)
        self.check_respread_file_exist = QtGui.QCheckBox()
        lay2.addWidget(self.check_respread_file_exist, 2, 1)
        
        lay4 = QtGui.QHBoxLayout()
        self.create_reapread = QtGui.QPushButton('', self)
        self.create_reapread.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_respread.table)))
        self.create_reapread.setIconSize(QtCore.QSize(24,24))
        self.connect(self.create_reapread, QtCore.SIGNAL('clicked()'), self.create_responsive_reading)
        
        self.create_reapread_sub = QtGui.QPushButton('', self)
        self.create_reapread_sub.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_respread_sub.table)))
        self.create_reapread_sub.setIconSize(QtCore.QSize(24,24))
        self.connect(self.create_reapread_sub, QtCore.SIGNAL('clicked()'), self.create_responsive_reading_subtitle)
        
        self.clear_reapread = QtGui.QPushButton('', self)
        self.clear_reapread.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_restore.table)))
        self.clear_reapread.setIconSize(QtCore.QSize(24,24))
        self.connect(self.clear_reapread, QtCore.SIGNAL('clicked()'), self.restore_responsive_reading_format)
        
        lay4.addWidget(self.create_reapread)
        lay4.addWidget(self.create_reapread_sub)
        lay4.addWidget(self.clear_reapread)
        
        #self.respread_edit = QtGui.QPlainTextEdit()
        #self.respread_edit.setFont(QtGui.QFont("Courier",9,True))
        #policy = self.sizePolicy()
        #policy.setVerticalStretch(1)
        #self.respread_edit.setSizePolicy(policy)
        
        layout.addRow(lay1)
        layout.addWidget(self.respread_format_tbl)
        layout.addRow(lay2)
        layout.addRow(lay4)
        self.respread_tab.setLayout(layout)
    
    def get_respread_path(self):
        startingDir = os.getcwd() 
        path = QtGui.QFileDialog.getExistingDirectory(None, 'Save folder', startingDir, 
        QtGui.QFileDialog.ShowDirsOnly)
        if not path: return
        self.respread_dest_path.setText(path)
        
    def restore_responsive_reading_format(self):
        #self.respread_edit.clear()
        #responsive_reading_format_key
        self.respread_format_tbl.setRowCount(len(const._responsive_reading_default_format))
        for ii, jj in enumerate(const._responsive_reading_default_format):
            item = QtGui.QTableWidgetItem(jj)
            item.setFlags(QtCore.Qt.ItemIsEnabled)
            fmt = const._responsive_reading_default_format[ii].split('|')
            self.respread_format_tbl.setItem(ii, 0, QtGui.QTableWidgetItem(fmt[0]))
            self.respread_format_tbl.setItem(ii, 1, QtGui.QTableWidgetItem(fmt[1]))
        
    def set_respread_slide_size(self):
        w,h = get_slide_size(self.choose_respread_slide_size.currentText())
        self.respread_format_tbl.item(0,1).setText("%2.3f, %2.3f"%(w,h))

    def check_respread_num_range(self, low, high):
        return False if low > high or low < 1 or high > 137 else True
        
    def create_responsive_reading_subtitle(self):
        import hymal
        self.global_message.appendPlainText('... Create RespRead Sub')
        try: 
            num_text = self.respread_num.text()
            if num_text.find('-') > 0:
                num_range = num_text.split('-')
                low, high = int(num_range[0]), int(num_range[1])
            else:
                low, high = int(num_text), int(num_text)
        except Exception as e: 
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid number', str(e),  QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('Invalid number: %s'%num_text)
            return
                
        if not self.check_respread_num_range(low, high):
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', 'Invalid number range %s'%num_text,  QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('Invalid number range: %s'%num_text)
            return
        
        skip_file_check = self.check_respread_file_exist.isChecked()
        
        for ir in range(low, high+1):
            rfn = '교독문 %03d-sub.pptx'%ir
            self.global_message.appendPlainText("--> %s"%rfn)
            sfn = os.path.join(self.respread_dest_path.text(), rfn)
            if skip_file_check and os.path.isfile(sfn):
                ans = QtGui.QMessageBox.question(self, 'Continue?', 
                        '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
                if ans == QtGui.QMessageBox.No: return
                else:
                    try:
                        os.remove(sfn)
                    except OSError as e:
                        e_str = str(e)
                        if func.access_denied(e_str):
                            e_str += "%s is already opened!"%sfn
                            msgcom.message_box(msgcom.message_error, e_str)
                            self.global_message.appendPlainText(e_str)
                            return
                try:
                    os.remove(sfn)
                except OSError as e:
                    e_str = str(e)
                    if func.access_denied(e_str):
                        e_str += "%s is already opened!"%sfn
                        msgcom.message_box(msgcom.message_error, e_str)
                        self.global_message.appendPlainText(e_str)
                        return
                
            _, rtext = hymal.get_responsive_reading_by_chapter(ir, self.hymnal_db_file_path.text())
            
            rsize = self.respread_format_tbl.item(0,1).text().split(',')
            rtext = rtext[0]
            wid, hgt = float(rsize[0]), float(rsize[1])
            pbx = self.ppt_respread_slide.textbox
            #dest_ppt = pptx.Presentation(_default_pptx_template)
            dest_ppt = pptx.Presentation(self.get_default_pptx())
            dest_ppt.slide_width = pptx.util.Inches(wid)
            dest_ppt.slide_height = pptx.util.Inches(hgt)
            blank_slide_layout = dest_ppt.slide_layouts[6]
                            
            for i, rt in enumerate(rtext):
                dest_slide = self.add_empty_slide(dest_ppt, blank_slide_layout, _color_white)
                txt_box = self.add_textbox(dest_slide, self.ppt_respread_slide, 0, hgt-pbx.hgt, wid, pbx.hgt)
                txt_f = txt_box.text_frame
                self.set_textbox(txt_f, 
                                MSO_AUTO_SIZE.NONE, 
                                const.get_textframe_vanchortype(self.ppt_respread_slide.textbox.vanchor),
                                const.get_paragraph_aligntype(self.ppt_respread_slide.textbox.pp_align),
                                #MSO_ANCHOR.MIDDLE, 
                                #MSO_ANCHOR.MIDDLE,
                                pbx.left_margin, pbx.top_margin,
                                pbx.right_margin, pbx.bottom_margin)
                self.add_paragraph(txt_f, re.sub('\(.+\)', '', rt).strip(), True)
            
            try:
                dest_ppt.save(sfn)
            except Exception as e:
                QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', str(e), 
                QtGui.QMessageBox.Yes)
                self.global_message.appendPlainText('... Error: %s'%str(e))
                dest_ppt = None
                return
                
            #ft = self.ppt_respread_slide.textbox.fill.fill_type
            #if self.ppt_respread_slide.textbox.fill.show and ft is _TEXTBOX_GRADIENT_FILL:
            #    self.fill_textbox(sfn, bool(ft))

        QtGui.QMessageBox.question(QtGui.QWidget(), 'Success', sfn, QtGui.QMessageBox.Yes)
        self.global_message.appendPlainText('... Create RespRead Sub: success\n')
        
    def set_respread_ppt_textbox(self, dest_slide, fmt):
        w,h = get_slide_size(self.choose_respread_slide_size.currentText())
        fmt_list = fmt.split(',')
        sx = float(fmt_list[0])
        sy = float(fmt_list[1])
        bkc = func.get_rgb(fmt_list[4].replace(':',','))
        txt_box = self.add_textbox(dest_slide, self.ppt_respread_slide, sx, sy, w-sx, h-sy)
        txt_box.fill.solid()
        txt_box.fill.fore_color.rgb = pptx.dml.color.RGBColor(bkc.r, bkc.g, bkc.b)
        txt_f = txt_box.text_frame
        self.set_textbox(txt_f, MSO_AUTO_SIZE.NONE, 
            #MSO_ANCHOR.TOP, 
            #MSO_ANCHOR.MIDDLE, 
            const.get_textframe_vanchortype(self.ppt_respread_slide.textbox.vanchor),
            const.get_paragraph_aligntype(self.ppt_respread_slide.textbox.pp_align),
            0, 0, 0, 0, True)
        return txt_box
            
    def set_respread_ppt_text(self, txt_frame, txt, fmt):
        fmt_list = fmt.split(',')
        p = txt_frame.add_paragraph()
        p.text = txt
        self.set_paragraph(p, 
                    const._responsive_reading_text_align[fmt_list[4].strip()], 
                    fmt_list[0], 
                    float(fmt_list[1]),
                    func.get_rgb(fmt_list[3].replace(':',',')),
                    bool(int(fmt_list[2])))
    
    def set_respread_ppt_newline(self, txt_frame, fmt):
        fmt_list = fmt.split(',')
        p = txt_frame.add_paragraph()
        p.text = ''
        self.set_paragraph(p, 
        const._responsive_reading_text_align['Left'], 
        fmt_list[0], float(fmt_list[1]), color._color_white, False)
        
    def create_responsive_reading(self):
        import hymal
        self.global_message.appendPlainText('... Create RespRead')
        
        try: 
            num_text = self.respread_num.text()
            if num_text.find('-') > 0:
                num_range = num_text.split('-')
                low, high = int(num_range[0]), int(num_range[1])
            else:
                low, high = int(num_text), int(num_text)
        except Exception as e: 
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid number', str(e),  QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('Invalid number: %s'%num_text)
            return
                
        if not self.check_respread_num_range(low, high):
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', 'Invalid number range %s'%num_text,  QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('Invalid number range: %s'%num_text)
            return
                     
        rfmt_list = []
        nfmt = self.respread_format_tbl.rowCount()        
        for i in range(0, nfmt):
            fmt_key = self.respread_format_tbl.item(i,0).text()
            fmt_str = self.respread_format_tbl.item(i,1).text()
            rfmt_list.append([fmt_key, fmt_str])
    
        skip_file_check = self.check_respread_file_exist.isChecked()
        
        for ir in range(low, high+1):
            rfn = '교독문 %03d.pptx'%ir
            self.global_message.appendPlainText("--> %s"%rfn)
            sfn = os.path.join(self.respread_dest_path.text(), rfn)
            if skip_file_check and os.path.isfile(sfn):
                ans = QtGui.QMessageBox.question(self, 'Continue?', 
                        '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
                if ans == QtGui.QMessageBox.No: return
                else:
                    try:
                        os.remove(sfn)
                    except OSError as e:
                        e_str = str(e)
                        if func.access_denied(e_str):
                            e_str += "%s is already opened!"%sfn
                            msgcom.message_box(msgcom.message_error, e_str)
                            self.global_message.appendPlainText(e_str)
                            return
            else: 
                try:
                    os.remove(sfn)
                except OSError as e:
                    e_str = str(e)
                    if func.access_denied(e_str):
                        e_str += "%s is already opened!"%sfn
                        msgcom.message_box(msgcom.message_error, e_str)
                        self.global_message.appendPlainText(e_str)
                        return

            rtitle, rtext = hymal.get_responsive_reading_by_chapter(ir, self.hymnal_db_file_path.text())
            
            if isinstance(rtitle, int) and rtitle < 0:
                self.global_message.appendPlainText(rtext)
                QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid DB', rtext, QtGui.QMessageBox.Yes)
                return            
            
            rsize = rfmt_list[0][1].split(',')
            
            #dest_ppt = pptx.Presentation(_default_pptx_template)
            dest_ppt = pptx.Presentation(self.get_default_pptx())
            dest_ppt.slide_width = pptx.util.Inches(float(rsize[0]))
            dest_ppt.slide_height = pptx.util.Inches(float(rsize[1]))
            blank_slide_layout = dest_ppt.slide_layouts[6]
            
            rtext = rtext[0]
            leftover = len(rtext)%2
            nread = int(len(rtext)/2)
            nslide = nread + leftover
            
            rfmt_dic = {
                const._responsive_reading_format_key[1]: self.set_respread_ppt_textbox,
                const._responsive_reading_format_key[2]: self.set_respread_ppt_text,
                const._responsive_reading_format_key[3]: self.set_respread_ppt_newline
            }
    
            for i in range(nslide):
                dest_slide = self.add_empty_slide(dest_ppt, blank_slide_layout, self.ppt_hymal.back_col)
                txt_box = self.set_respread_ppt_textbox(dest_slide, rfmt_list[1][1])
                txt_f = txt_box.text_frame
                txt_f.margin_left   = pptx.util.Inches(0.1 )
                txt_f.margin_top    = pptx.util.Inches(0.05 )
                txt_f.margin_right  = pptx.util.Inches(0.1)
                txt_f.margin_bottom = pptx.util.Inches(0.05)
                
                if i is 0:
                    self.set_respread_ppt_text(txt_f, '교독문 %d번 (%s)'%(ir,rtitle), rfmt_list[2][1]) 
                    self.set_respread_ppt_newline(txt_f, rfmt_list[3][1])
    
                if leftover and i is nslide-1:
                    rt1 = re.sub('\(.+\)', '', rtext.pop(0)).strip()
                    self.set_respread_ppt_text(txt_f, '(다같이)', rfmt_list[7][1]) 
                    self.set_respread_ppt_text(txt_f, rt1, rfmt_list[8][1]) 
                else:
                    rt1 = rtext.pop(0).strip()
                    self.set_respread_ppt_text(txt_f, '(사회자)', rfmt_list[4][1]) 
                    self.set_respread_ppt_text(txt_f, rt1, rfmt_list[5][1]) 
                    self.set_respread_ppt_newline(txt_f, rfmt_list[6][1])
                    rt2 = rtext.pop(0).strip()
                    self.set_respread_ppt_text(txt_f, '(다같이)', rfmt_list[7][1]) 
                    self.set_respread_ppt_text(txt_f, rt2, rfmt_list[8][1])                
            try:
                dest_ppt.save(sfn)
            except Exception as e:
                QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', str(e), 
                QtGui.QMessageBox.Yes)
                self.global_message.appendPlainText('... Error: %s'%str(e))
                dest_ppt = None
                return
            
        QtGui.QMessageBox.question(QtGui.QWidget(), 'Success', sfn, 
            QtGui.QMessageBox.Yes)
        self.global_message.appendPlainText('... Create RespRead: success\n')
        
    def fx_tab_UI(self):
        import msoDash
        layout = QtGui.QFormLayout()
        ohlyout = QtGui.QHBoxLayout()
        ohlyout.addWidget(QtGui.QLabel("OUTLINE EFFECT"))
        self.fx_show_outline = QtGui.QCheckBox()
        self.fx_show_outline.setChecked(self.ppt_subtitle_slide.textbox.fx.show_outline)
        ohlyout.addWidget(self.fx_show_outline)

        self.fx_outline_tbl = QtGui.QTableWidget()
        font = QtGui.QFont("Arial",9,True)
        self.fx_outline_tbl.setFont(font)
        self.fx_outline_tbl.horizontalHeader().hide()
        self.fx_outline_tbl.verticalHeader().hide()
        self.fx_outline_tbl.setColumnCount(3)
        self.fx_outline_tbl.setHorizontalHeaderItem(0, QtGui.QTableWidgetItem("ddd"))
        self.fx_outline_tbl.setHorizontalHeaderItem(1, QtGui.QTableWidgetItem("Data"))
        self.fx_outline_tbl.setHorizontalHeaderItem(2, QtGui.QTableWidgetItem("aaa"))

        header = self.fx_outline_tbl.horizontalHeader()
        header.setResizeMode(0, QtGui.QHeaderView.ResizeToContents)
        header.setResizeMode(1, QtGui.QHeaderView.Stretch)
        header.setResizeMode(2, QtGui.QHeaderView.ResizeToContents)
        
        info = ["Show", "Color", "Line", "Dash", "Transp", "weight"]
        self.fx_outline_tbl.setRowCount(len(info))
        for ii, jj in enumerate(info):
            item = QtGui.QTableWidgetItem(jj)
            item.setFlags(QtCore.Qt.ItemIsEnabled)
            self.fx_outline_tbl.setItem(ii, 0, item)
            self.fx_outline_tbl.setItem(ii, 1, QtGui.QTableWidgetItem(""))
            
        self.fx_outline_show = QtGui.QCheckBox()
        self.fx_outline_show.setChecked(True)
        self.fx_outline_tbl.setCellWidget(0, 1, self.fx_outline_show)

        self.fx_outline_col_picker = QtGui.QPushButton('', self)
        self.fx_outline_col_picker.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_color_picker.table)))
        self.fx_outline_col_picker.setIconSize(QtCore.QSize(16,16))
        self.connect(self.fx_outline_col_picker, QtCore.SIGNAL('clicked()'), self.pick_text_outline_color)
        self.fx_outline_tbl.setCellWidget(1,2, self.fx_outline_col_picker)
        self.fx_outline_tbl.item(1,1).setText(str(self.ppt_subtitle_slide.textbox.fx.outline.col))

        style = msoLine.get_linestyle_list()
        self.fx_outline_style = QtGui.QComboBox()
        self.fx_outline_style.addItems(style)
        self.fx_outline_tbl.setCellWidget(2,1, self.fx_outline_style)

        dash = msoDash.get_dashstyle_list()
        dash.insert(0, "N/A")
        self.fx_outline_dash_style = QtGui.QComboBox()
        self.fx_outline_dash_style.addItems(dash)
        self.fx_outline_dash_style.setCurrentIndex(1) # solid 
        self.fx_outline_tbl.setCellWidget(3,1, self.fx_outline_dash_style)

        self.fx_outline_tbl.item(4,1).setText("%f"%self.ppt_subtitle_slide.textbox.fx.outline.transprancy)
        self.fx_outline_tbl.item(5,1).setText("%f"%self.ppt_subtitle_slide.textbox.fx.outline.weight)

        self.fx_outline_tbl.resizeRowsToContents()

        shlyout = QtGui.QHBoxLayout()
        shlyout.addWidget(QtGui.QLabel("SHADOW EFFECT"))
        self.fx_show_shadow = QtGui.QCheckBox()
        self.fx_show_shadow.setChecked(self.ppt_subtitle_slide.textbox.fx.show_shadow)
        shlyout.addWidget(self.fx_show_shadow)

        self.fx_shadow_tbl = QtGui.QTableWidget()
        self.fx_shadow_tbl.resizeRowsToContents()	
        self.fx_shadow_tbl.setFont(font)
        self.fx_shadow_tbl.horizontalHeader().hide()
        self.fx_shadow_tbl.verticalHeader().hide()
        self.fx_shadow_tbl.setColumnCount(3)
        self.fx_shadow_tbl.setHorizontalHeaderItem(0, QtGui.QTableWidgetItem("ddd"))
        self.fx_shadow_tbl.setHorizontalHeaderItem(1, QtGui.QTableWidgetItem("Data"))
        self.fx_shadow_tbl.setHorizontalHeaderItem(2, QtGui.QTableWidgetItem("aaa"))

        info = ["Style", "Offset(X)", "Offset(Y)", "Blur", "Transp"]
        self.fx_shadow_tbl.setRowCount(len(info))
        for ii, jj in enumerate(info):
            item = QtGui.QTableWidgetItem(jj)
            item.setFlags(QtCore.Qt.ItemIsEnabled)
            self.fx_shadow_tbl.setItem(ii, 0, item)
            self.fx_shadow_tbl.setItem(ii, 1, QtGui.QTableWidgetItem(""))

        self.fx_shadow_style = QtGui.QComboBox(self)
        s_style = ["Inner", "Outer", "Mixed"]
        self.fx_shadow_style.addItems(s_style)
        self.fx_shadow_style.setCurrentIndex(self.ppt_subtitle_slide.textbox.fx.shadow.Style-1)
        self.fx_shadow_tbl.setCellWidget(0,1,self.fx_shadow_style)

        self.fx_shadow_tbl.item(1,1).setText("%d"%self.ppt_subtitle_slide.textbox.fx.shadow.OffsetX)
        self.fx_shadow_tbl.item(2,1).setText("%d"%self.ppt_subtitle_slide.textbox.fx.shadow.OffsetY)
        self.fx_shadow_tbl.item(3,1).setText("%d"%self.ppt_subtitle_slide.textbox.fx.shadow.Blur)
        self.fx_shadow_tbl.item(4,1).setText("%f"%self.ppt_subtitle_slide.textbox.fx.shadow.Transparency)

        self.fx_shadow_tbl.resizeRowsToContents()

        layout.addRow(ohlyout)
        layout.addRow(self.fx_outline_tbl)
        layout.addRow(shlyout)
        layout.addRow(self.fx_shadow_tbl)
        self.fx_tab.setLayout(layout)
        self.global_message.appendPlainText("... Fx Tab UI created\n")

    def clear_global_message(self):
        self.global_message.clear()

    def message_tab_UI(self):
        layout = QtGui.QVBoxLayout()

        clear_btn = QtGui.QPushButton('Clear', self)
        clear_btn.clicked.connect(self.clear_global_message)

        self.global_message = QtGui.QPlainTextEdit()
        self.global_message.setFont(QtGui.QFont("Courier",9,True))
        policy = self.sizePolicy()
        policy.setVerticalStretch(1)
        self.global_message.setSizePolicy(policy)
        self.global_message.setFont(QtGui.QFont( "Courier,15,-1,2,50,0,0,0,1,0"))
        layout.addWidget(clear_btn)
        layout.addWidget(self.global_message)
        self.message_tab.setLayout(layout)
        self.global_message.appendPlainText("... Message Tab UI Created")

    def clear_txtppt(self):
        self.txtppt_edit.clear()

    def change_txtppt_save_path(self):
        startingDir = os.getcwd() 
        path = QtGui.QFileDialog.getExistingDirectory(None, 'Save folder', startingDir, QtGui.QFileDialog.ShowDirsOnly)
        if not path: return
        self.txtppt_save_path.setText(path)

    def convert_txt_pptx(self):
        text = self.txtppt_edit.toPlainText().split('\n')

        for i, t in enumerate(text):
            text[i] = t.strip()

        size = len(text) 
        #print(size)
        if size == 1:
            line_text = text
        else:
            #https://www.geeksforgeeks.org/python-split-list-into-lists-by-particular-value/
            idx_list = [idx + 1 for idx, val in enumerate(text) if val == ''] 
            line_text = [text[i:j] 
                            for i, j in zip([0] + idx_list, 
                                idx_list + (
                                    [size] if idx_list[-1] != size else []
                                    )
                                )
                        ] 

        fn = self.txtppt_fname.text()
        if fn.endswith('.pptx'):
            sfn = os.path.join(self.txtppt_save_path.text(), fn)
        else:
            sfn = os.path.join(self.txtppt_save_path.text(), '%s.pptx'%fn)

        if os.path.isfile(sfn):
            ans = QtGui.QMessageBox.question(self, 'Continue?', 
                    '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
            if ans == QtGui.QMessageBox.No: return
            else:
                try:
                    os.remove(sfn)
                except OSError as e:
                    if func.access_denied(e_str):
                        e_str += "%s is already opened!"%sfn
                        msgcom.message_box(msgcom.message_error, e_str)
                        return

        self.global_message.appendPlainText('... Create TxtPPT')
        #dest_ppt = pptx.Presentation(_default_pptx_template)
        dest_ppt = pptx.Presentation(self.get_default_pptx())
        dest_ppt.slide_width = pptx.util.Inches(self.ppt_txtppt_slide.slide.wid)
        dest_ppt.slide_height = pptx.util.Inches(self.ppt_txtppt_slide.slide.hgt)
        blank_slide_layout = dest_ppt.slide_layouts[6]

        self.global_message.appendPlainText('Total slide: %d'%len(line_text))
        self.global_message.appendPlainText('%s'%str(self.ppt_txtppt_slide))

        sx  = self.ppt_txtppt_slide.textbox.sx 
        sy  = self.ppt_txtppt_slide.textbox.sy 
        wid = self.ppt_txtppt_slide.textbox.wid
        hgt = self.ppt_txtppt_slide.textbox.hgt
        font_name = self.ppt_txtppt_slide.textbox.font_name
        font_col  = self.ppt_txtppt_slide.textbox.font_col
        font_size = self.ppt_txtppt_slide.textbox.font_size
        back_col  = self.ppt_txtppt_slide.slide.back_col
        
        for in_text in line_text:
            dest_slide = self.add_empty_slide(dest_ppt, blank_slide_layout, back_col)
            txt_box = self.add_textbox(dest_slide, self.ppt_txtppt_slide, sx, sy, wid, hgt)
            txt_f = txt_box.text_frame
            self.set_textbox(txt_f, 
                MSO_AUTO_SIZE.NONE, 
                #MSO_ANCHOR.MIDDLE, 
                #MSO_ANCHOR.MIDDLE, 
                const.get_textframe_vanchortype(self.ppt_txtppt_slide.textbox.vanchor),
                const.get_paragraph_aligntype(self.ppt_txtppt_slide.textbox.pp_align),
                0, 0, 0, 0)

            for l2 in in_text:
                if l2 == '': continue
                p = txt_f.add_paragraph()
                p.text = l2
                self.set_paragraph(p, 
                    #PP_ALIGN.CENTER, 
                    const.get_paragraph_aligntype(self.ppt_txtppt_slide.textbox.pp_align),
                    font_name, font_size, font_col, True)
        try:
            dest_ppt.save(sfn)
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Error: %s'%e)
            dest_ppt = None
            return

        self.global_message.appendPlainText('... Create TxtPPT: success\n')
        QtGui.QMessageBox.question(QtGui.QWidget(), 'Success', sfn, QtGui.QMessageBox.Yes)

    def txtppt_tab_UI(self):
        layout = QtGui.QFormLayout()
        laybtn = QtGui.QHBoxLayout()

        clear_btn = QtGui.QPushButton('Clear', self)
        clear_btn.clicked.connect(self.clear_txtppt)

        run_btn = QtGui.QPushButton('Create', self)
        run_btn.clicked.connect(self.convert_txt_pptx)
        laybtn.addWidget(clear_btn)
        laybtn.addWidget(run_btn)

        laypub = QtGui.QGridLayout()
        #layfolder = QtGui.QHBoxLayout()
        #layfolder.addWidget(QtGui.QLabel('Dest'))
        laypub.addWidget(QtGui.QLabel('Dest'), 0, 0)
        self.txtppt_save_path  = QtGui.QLineEdit(os.getcwd())
        self.txtppt_save_path_btn = QtGui.QPushButton('', self)
        self.txtppt_save_path_btn.clicked.connect(self.change_txtppt_save_path)
        self.txtppt_save_path_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.txtppt_save_path_btn.setIconSize(QtCore.QSize(16,16))
        #layfolder.addWidget(self.txtppt_save_path)
        #layfolder.addWidget(self.txtppt_save_path_btn)
        laypub.addWidget(self.txtppt_save_path, 0, 1)
        laypub.addWidget(self.txtppt_save_path_btn, 0, 2)
        
        #layfile = QtGui.QHBoxLayout()
        #layfile.addWidget(QtGui.QLabel('File'))
        laypub.addWidget(QtGui.QLabel('File'), 1, 0)
        self.txtppt_fname = QtGui.QLineEdit('txtppt.pptx')
        #layfile.addWidget(self.txtppt_fname)
        
        self.txtppt_slide_data_btn = QtGui.QPushButton('', self)
        self.txtppt_slide_data_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_slide_setting.table)))
        self.txtppt_slide_data_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.txtppt_slide_data_btn, QtCore.SIGNAL('clicked()'), self.set_txtppt_slide_data)
         
        laypub.addWidget(self.txtppt_fname, 1, 1)
        laypub.addWidget(self.txtppt_slide_data_btn, 1, 2)
        
        self.txtppt_edit = QtGui.QPlainTextEdit()
        self.txtppt_edit.setFont(QtGui.QFont("Courier",9,True))
        policy = self.sizePolicy()
        policy.setVerticalStretch(1)
        self.txtppt_edit.setSizePolicy(policy)
        self.txtppt_edit.setFont(QtGui.QFont( "Courier,15,-1,2,50,0,0,0,1,0"))
        
        layout.addRow(laybtn)
        #layout.addRow(layfolder)
        #layout.addRow(layfile)
        layout.addRow(laypub)
        layout.addWidget(self.txtppt_edit)
        self.txtppt_tab.setLayout(layout)
        
    def set_txtppt_slide_data(self):
        res = slide_data_dlg.get_slide_data(self.ppt_txtppt_slide, sub_title=False)
        if res == 0: return
        
    def hymal_tab_UI(self):
        layout = QtGui.QFormLayout()
        
        lay = QtGui.QHBoxLayout()
        lay.addWidget(QtGui.QLabel("Hymnal Chap"))
        self.hymal_num = QtGui.QLineEdit()
        self.hymal_num.setText("%d"%self.ppt_hymal_slide.chap)
        self.hymal_num.setToolTip("1-639")
        lay.addWidget(self.hymal_num)
        layout.addRow(lay)
        
        self.hymal_info_table = QtGui.QTableWidget()
        #font = QtGui.QFont("Arial",9,True)
        #self.hymal_info_table.setFont(font)
        self.hymal_info_table.horizontalHeader().hide()
        self.hymal_info_table.verticalHeader().hide()
        self.hymal_info_table.setColumnCount(3)
        self.hymal_info_table.setHorizontalHeaderItem(0, QtGui.QTableWidgetItem("ddd"))
        self.hymal_info_table.setHorizontalHeaderItem(1, QtGui.QTableWidgetItem("Data"))
        self.hymal_info_table.setHorizontalHeaderItem(2, QtGui.QTableWidgetItem("aaa"))
        
        header = self.hymal_info_table.horizontalHeader()
        header.setResizeMode(0, QtGui.QHeaderView.ResizeToContents)
        header.setResizeMode(1, QtGui.QHeaderView.Stretch)
        header.setResizeMode(2, QtGui.QHeaderView.ResizeToContents)
        
        info = [#"Chapter", 
                "Textbox(sx)", "Textbox(sy)", "Textbox(wid)",
                "Textbox(hgt)", "Font", "Font(RGB)", "Size(pt)", "Back(RGB)"]
    
        self.hymal_info_table.setRowCount(len(info))
                
        for ii, jj in enumerate(info):
            item = QtGui.QTableWidgetItem(jj)
            item.setFlags(QtCore.Qt.ItemIsEnabled)
            self.hymal_info_table.setItem(ii, 0, item)
            self.hymal_info_table.setItem(ii, 1, QtGui.QTableWidgetItem(""))
        
        self.hymal_textbox_font_picker = QtGui.QPushButton('', self)
        self.hymal_textbox_font_picker.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_font_picker.table)))
        self.hymal_textbox_font_picker.setIconSize(QtCore.QSize(16,16))
        self.connect(self.hymal_textbox_font_picker, QtCore.SIGNAL('clicked()'), self.pick_hymal_font)
        self.hymal_info_table.setCellWidget(4,2, self.hymal_textbox_font_picker)
        
        self.hymal_font_col_picker = QtGui.QPushButton('', self)
        self.hymal_font_col_picker.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_color_picker.table)))
        self.hymal_font_col_picker.setIconSize(QtCore.QSize(16,16))
        self.connect(self.hymal_font_col_picker, QtCore.SIGNAL('clicked()'), self.pick_hymal_font_color)
        self.hymal_info_table.setCellWidget(5,2, self.hymal_font_col_picker)
    
        self.hymal_back_col_picker = QtGui.QPushButton('', self)
        self.hymal_back_col_picker.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_color_picker.table)))
        self.hymal_back_col_picker.setIconSize(QtCore.QSize(16,16))
        self.connect(self.hymal_back_col_picker, QtCore.SIGNAL('clicked()'), self.pick_hymal_bk_color)
        self.hymal_info_table.setCellWidget(7,2, self.hymal_back_col_picker)
        
        #self.hymal_info_table.item(0,1).setText("%d"%self.ppt_hymal_slide.chap)
        self.hymal_info_table.item(0,1).setText("%f"%float(self.ppt_hymal_slide.textbox.sx ))
        self.hymal_info_table.item(1,1).setText("%f"%float(self.ppt_hymal_slide.textbox.sy ))
        self.hymal_info_table.item(2,1).setText("%f"%float(self.ppt_hymal_slide.textbox.wid))
        self.hymal_info_table.item(3,1).setText("%f"%float(self.ppt_hymal_slide.textbox.hgt))
        self.hymal_info_table.item(4,1).setText(self.ppt_hymal_slide.textbox.font_name)
        c = self.ppt_hymal_slide.textbox.font_col
        c1="%03d,%03d,%03d"%(c.r,c.g,c.b)
        self.hymal_info_table.item(5,1).setText(c1)
        self.hymal_info_table.item(6,1).setText("%d"%self.ppt_hymal_slide.textbox.font_size)
        c = self.ppt_hymal_slide.slide.back_col
        c1="%03d,%03d,%03d"%(c.r,c.g,c.b)
        self.hymal_info_table.item(7,1).setText(c1)
        self.hymal_info_table.resizeRowsToContents()			
    
        db_layout = QtGui.QGridLayout()
        db_layout.addWidget(QtGui.QLabel('Slide Size'), 0, 0)
        self.choose_hymal_slide_size = QtGui.QComboBox(self)
        self.choose_hymal_slide_size.addItems(const._slide_size_type)
        self.choose_hymal_slide_size.setCurrentIndex(self.ppt_hymal_slide.slide.size_index)
        self.choose_hymal_slide_size.currentIndexChanged.connect(self.set_hymal_slide_size)
        db_layout.addWidget(self.choose_hymal_slide_size, 0, 1)
        
        self.hymal_slide_data_btn = QtGui.QPushButton('', self)
        self.hymal_slide_data_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_slide_setting.table)))
        self.hymal_slide_data_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.hymal_slide_data_btn, QtCore.SIGNAL('clicked()'), self.set_hymal_slide_data)
        db_layout.addWidget(self.hymal_slide_data_btn, 0, 2)

        db_layout.addWidget(QtGui.QLabel('Dest'), 1, 0)
        #self.ppt_hymal_slide.save_path = os.getcwd()
        #self.ppt_hymal_slide.save_path = os.getcwd()
        self.hymal_save_path  = QtGui.QLineEdit(os.getcwd())#self.ppt_hymal_slide.save_path)
        self.hymal_save_path_btn = QtGui.QPushButton('', self)
        self.hymal_save_path_btn.clicked.connect(self.change_hymal_save_path)
        self.hymal_save_path_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.hymal_save_path_btn.setIconSize(QtCore.QSize(16,16))
        db_layout.addWidget(self.hymal_save_path, 1, 1)
        db_layout.addWidget(self.hymal_save_path_btn, 1, 2)
        
        db_layout.addWidget(QtGui.QLabel('Hymal DB'), 2, 0)
        self.hymnal_db_file_path  = QtGui.QLineEdit(hymal._default_db_file)
        self.hymnal_db_file_path_btn = QtGui.QPushButton('', self)
        self.hymnal_db_file_path_btn.clicked.connect(self.change_hymal_db_file)
        self.hymnal_db_file_path_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.hymnal_db_file_path_btn.setIconSize(QtCore.QSize(16,16))
        db_layout.addWidget(self.hymnal_db_file_path, 2, 1)
        db_layout.addWidget(self.hymnal_db_file_path_btn, 2, 2)
        
        db_layout.addWidget(QtGui.QLabel('File Check'), 3, 0)
        self.hymnal_file_exist_chk = QtGui.QCheckBox()
        self.hymnal_file_exist_chk.setChecked(False)
        db_layout.addWidget(self.hymnal_file_exist_chk, 3,1)
        
        
        self.hymal_fit_textbox_slide_btn = QtGui.QPushButton('', self)
        self.hymal_fit_textbox_slide_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_fit.table)))
        self.hymal_fit_textbox_slide_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.hymal_fit_textbox_slide_btn, QtCore.SIGNAL('clicked()'), self.fit_textbox_to_slide)
        db_layout.addWidget(self.hymal_fit_textbox_slide_btn, 3, 2)
        
        
        run_layout = QtGui.QHBoxLayout()
        self.hymal_convert_bth = QtGui.QPushButton('', self)
        self.hymal_convert_bth.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_convert.table)))
        self.hymal_convert_bth.setIconSize(QtCore.QSize(24,24))
        self.hymal_convert_bth.clicked.connect(self.create_hymal_ppt)
        run_layout.addWidget(self.hymal_convert_bth)
            
        layout.addRow(self.hymal_info_table)
        #layout.addRow(size_layout)
        layout.addRow(db_layout)
        layout.addRow(run_layout)
        self.hymal_tab.setLayout(layout)
        self.global_message.appendPlainText('... Hymal tab UI created')
        
    def fit_textbox_to_slide(self):
        self.ppt_hymal_slide.textbox.sx = 0
        self.ppt_hymal_slide.textbox.sy = 0 
        self.ppt_hymal_slide.textbox.wid = self.ppt_hymal_slide.slide.wid
        self.ppt_hymal_slide.textbox.hgt = self.ppt_hymal_slide.slide.hgt
        #self.update_hymal_slide_data()
        
    def update_hymal_slide_info(self):
        self.hymal_info_table.item(0,1).setText("%f"%float(self.ppt_hymal_slide.textbox.sx ))
        self.hymal_info_table.item(1,1).setText("%f"%float(self.ppt_hymal_slide.textbox.sy ))
        self.hymal_info_table.item(2,1).setText("%f"%float(self.ppt_hymal_slide.textbox.wid))
        self.hymal_info_table.item(3,1).setText("%f"%float(self.ppt_hymal_slide.textbox.hgt))
        self.hymal_info_table.item(4,1).setText(self.ppt_hymal_slide.textbox.font_name)
        c = self.ppt_hymal_slide.textbox.font_col
        c1="%03d,%03d,%03d"%(c.r,c.g,c.b)
        self.hymal_info_table.item(5,1).setText(c1)
        self.hymal_info_table.item(6,1).setText("%d"%self.ppt_hymal_slide.textbox.font_size)
        c = self.ppt_hymal_slide.slide.back_col
        c1="%03d,%03d,%03d"%(c.r,c.g,c.b)
        self.hymal_info_table.item(7,1).setText(c1)

    def update_hymal_slide_data(self):
        self.ppt_hymal_slide.textbox.sx  = float(self.hymal_info_table.item(0,1).text())
        self.ppt_hymal_slide.textbox.sy  = float(self.hymal_info_table.item(1,1).text())
        self.ppt_hymal_slide.textbox.wid = float(self.hymal_info_table.item(2,1).text())
        self.ppt_hymal_slide.textbox.hgt = float(self.hymal_info_table.item(3,1).text())
        self.ppt_hymal_slide.textbox.font_name = self.hymal_info_table.item(4,1).text()
        self.ppt_hymal_slide.textbox.font_col  = func.get_rgb(self.hymal_info_table.item(5,1).text())
        self.ppt_hymal_slide.textbox.font_size = float(self.hymal_info_table.item(6,1).text())
        self.ppt_hymal_slide.slide.back_col    = func.get_rgb(self.hymal_info_table.item(7,1).text())
    
    def set_hymal_slide_data(self):
        self.update_hymal_slide_data()
        res = slide_data_dlg.get_slide_data(self.ppt_hymal_slide, sub_title=False)
        if res == 0: return
        
        new_index = self.ppt_hymal_slide.slide.size_index
        cur_index = self.choose_hymal_slide_size.currentIndex()
        if new_index != cur_index:
            self.choose_hymal_slide_size.setCurrentIndex(new_index)
            #w = self.ppt_hymal_slide.slide.wid
            #h = self.ppt_hymal_slide.slide.hgt
        self.update_hymal_slide_info()
            
    # direct chanage from combo box
    def set_hymal_slide_size(self):
        cur_index = self.choose_hymal_slide_size.currentIndex()
        new_index = self.ppt_hymal_slide.slide.size_index

        if cur_index != new_index:
            w,h = get_slide_size(self.choose_hymal_slide_size.currentText())
            self.ppt_hymal_slide.slide.wid = w
            self.ppt_hymal_slide.slide.hgt = h
            self.ppt_hymal_slide.slide.size_index = cur_index
            self.fit_textbox_to_slide()
            self.update_hymal_slide_info()
        
    def change_hymal_db_file(self):
        file = QtGui.QFileDialog.getOpenFileName(self, "Choose Hymal DB", directory=os.getcwd())
        if not file: return
        self.hymnal_db_file_path.setText(file) 
        
    def change_hymal_save_path(self):
        startingDir = os.getcwd() 
        path = QtGui.QFileDialog.getExistingDirectory(None, 'Save folder', startingDir, QtGui.QFileDialog.ShowDirsOnly)
        if not path: return
        self.hymal_save_path.setText(path)
    
    def pick_hymal_bk_color(self):
        bc = self.ppt_hymal_slide.slide.back_col
        col = QtGui.QColorDialog.getColor(QtGui.QColor(bc.r, bc.g, bc.b))
        if col.isValid():
            r,g,b,a = col.getRgb()
            self.ppt_hymal_slide.slide.back_col = colorppt_color(r,g,b)
            self.hymal_info_table.item(7,1).setText("%03d,%03d,%03d"%(r, g, b))
        
    def pick_hymal_font_color(self):
        fc = self.ppt_hymal_slide.textbox.font_col
        col = QtGui.QColorDialog.getColor(QtGui.QColor(fc.r, fc.g, fc.b))
        if col.isValid():
            r,g,b,a = col.getRgb()
            self.ppt_hymal_slide.textbox.font_col = color.ppt_color(r,g,b)
            self.hymal_info_table.item(5,1).setText("%03d,%03d,%03d"%(r, g, b))
        
    def pick_hymal_font(self):
        font, valid = QtGui.QFontDialog.getFont()
        if valid: 
            self.ppt_hymal_slide.textbox.font_name = font.family()
            self.hymal_info_table.item(4,1).setText(font.family())	
            
    # hymal: 1-639 
    def check_hymal_num_range(self, low, high):
        return False if low > high or low < 1 or high > 639 else True
        
    def create_hymal_ppt(self):
        import hymal
        import titnum
        
        self.global_message.appendPlainText("... Create Hymal")
        
        try: 
            num_text = self.hymal_num.text()
            if num_text.find('-') > 0:
                num_range = num_text.split('-')
                low, high = int(num_range[0]), int(num_range[1])
            else:
                low, high = int(num_text), int(num_text)
        except Exception as e: 
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid number', str(e),  QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('Invalid number: %s'%num_text)
            return
        
        if not self.check_hymal_num_range(low, high):
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', 'Invalid number range %s'%num_text,  QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('Invalid number range: %s'%num_text)
            return
          
        col_str = self.hymal_info_table.item(7,1).text()
        bk_col = func.get_rgb(col_str)
        if not bk_col:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid Back color', 
            "Comma separated {}".format(col_str), QtGui.QMessageBox.Yes)
            return
    
        col_str = self.hymal_info_table.item(5,1).text()
        ft_col = func.get_rgb(col_str)
        if not ft_col:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid Font color', 
            "Comma separated {}".format(col_str), QtGui.QMessageBox.Yes)
            return

        skip_file_check = self.hymnal_file_exist_chk.isChecked()
            
        for chap in range(low, high+1):
            title, lyric = hymal.get_hymal_by_chapter(chap, self.hymnal_db_file_path.text())
            if isinstance(title, int) and title < 0:
                QtGui.QMessageBox.question(QtGui.QWidget(), 'Invalid DB', lyric, QtGui.QMessageBox.Yes)
                self.global_message.appendPlainText(lyric)
                return
                
            title = re.sub("[<br><b></b>]", '', title)
            fname = "%03d-%s.pptx"%(chap, title)
            self.global_message.appendPlainText("Title: %s\nChap: %d"%(title, chap))
            sfn = os.path.join(self.hymal_save_path.text(), fname)
            
            if skip_file_check and os.path.isfile(sfn):
                ans = QtGui.QMessageBox.question(self, 'Continue?', 
                        '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
                if ans == QtGui.QMessageBox.No: return
                else:
                    try:
                        os.remove(sfn)
                    except OSError as e:
                        e_str = str(e)
                        if func.access_denied(e_str):
                            e_str += "%s is already opened!"%sfn
                            msgcom.message_box(msgcom.message_error, e_str)
                            return
            else: 
                try:
                    os.remove(sfn)
                except OSError as e:
                    e_str = str(e)
                    if func.access_denied(e_str):
                        e_str += "%s is already opened!"%sfn
                        msgcom.message_box(msgcom.message_error, e_str)
                        return
                    
            dest_ppt = pptx.Presentation(self.get_default_pptx())
            dest_ppt.slide_width = pptx.util.Inches(self.ppt_hymal_slide.slide.wid)
            dest_ppt.slide_height = pptx.util.Inches(self.ppt_hymal_slide.slide.hgt)
            blank_slide_layout = dest_ppt.slide_layouts[6]

            l_num = 1            
            for l_list in lyric:
                dest_slide = self.add_empty_slide(dest_ppt, blank_slide_layout, 
                self.ppt_hymal_slide.slide.back_col)
                txt_box = self.add_textbox(dest_slide, 
                                        self.ppt_hymal_slide,
                                        self.ppt_hymal_slide.textbox.sx,
                                        self.ppt_hymal_slide.textbox.sy,
                                        self.ppt_hymal_slide.textbox.wid,
                                        self.ppt_hymal_slide.textbox.hgt)
                txt_f = txt_box.text_frame
                self.set_textbox(txt_f, MSO_AUTO_SIZE.NONE, 
                const.get_textframe_vanchortype(self.ppt_hymal_slide.textbox.vanchor),
                const.get_paragraph_aligntype(self.ppt_hymal_slide.textbox.pp_align),
                0, 0, 0, 0)

                chorus = False
                for l_text in l_list:
                    p = txt_f.add_paragraph()
                    if l_text.find(hymal._corus_delimiter) >= 0:
                        l_text = l_text.replace(hymal._corus_delimiter, '')
                        chorus = True
                        
                    p.text = l_text
                    self.set_paragraph(p, 
                                const.get_paragraph_aligntype(self.ppt_hymal_slide.textbox.pp_align),
                                self.ppt_hymal_slide.textbox.font_name, 
                                self.ppt_hymal_slide.textbox.font_size,
                                self.ppt_hymal_slide.textbox.font_col, 
                                True)
                p = txt_f.add_paragraph()
                p.text ='\n'
                self.set_paragraph(p, 
                                const.get_paragraph_aligntype(self.ppt_hymal_slide.textbox.pp_align),
                                self.ppt_hymal_slide.textbox.font_name, 
                                10,
                                self.ppt_hymal_slide.textbox.font_col, 
                                True)
                                
                p = txt_f.add_paragraph()
                
                if len(lyric) == 1:
                    p.text = '(찬송가 %d장)'%(chap)
                elif chorus:
                    p.text = '(찬송가 %d장 후렴)'%(chap)
                else:
                    p.text = '(찬송가 %d장 %d절)'%(chap, l_num)
                    l_num += 1
                    
                self.set_paragraph(p, 
                                #PP_ALIGN.CENTER, 
                                const.get_paragraph_aligntype(self.ppt_hymal_slide.textbox.pp_align),
                                self.ppt_hymal_slide.textbox.font_name, 
                                20,
                                self.ppt_hymal_slide.textbox.font_col, 
                                True)
                    
            try:
                dest_ppt.save(sfn)
            except Exception as e:
                QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', str(e), 
                QtGui.QMessageBox.Yes)
                self.global_message.appendPlainText('... Error: %s'%str(e))
                dest_ppt = None
                return
            
        QtGui.QMessageBox.question(QtGui.QWidget(), 'Completed!', sfn, QtGui.QMessageBox.Yes)
        self.global_message.appendPlainText('... Create Hymal: success\n')
            
    def text_outline_state_changed(self):
        if self.text_outline.isChecked():
            self.ppt_subtitle_slide.textbox.fx.show_outline = True
        else:
            self.ppt_subtitle_slide.textbox.fx.show_outline = False
            
    def textbox_word_wrap_state_changed(self):
        if self.textbox_deep_copy.isChecked():
            self.ppt_subtitle_slide.textbox.word_wrap = True
        else:
            self.ppt_subtitle_slide.textbox.word_wrap = False
    
    def textbox_deep_copy_state_changed(self):
        if self.textbox_deep_copy.isChecked():
            self.ppt_subtitle_slide.slide.deep_copy = True
        else:
            self.ppt_subtitle_slide.slide.deep_copy = False
            
    def pick_text_outline_color(self):
        lc = self.ppt_subtitle_slide.textbox.fx.outline.col
        col = QtGui.QColorDialog.getColor(QtGui.QColor(lc.r, lc.g, lc.b))
        if col.isValid():
            r,g,b,a = col.getRgb()
            self.ppt_subtitle_slide.textbox.fx.outline.col = color.ppt_color(r,g,b)
            self.fx_outline_tbl.item(1,1).setText("%03d,%03d,%03d"%(r, g, b))
            
    def ppt_tab_UI(self):
        layout = QtGui.QFormLayout()
        open_layout  = QtGui.QHBoxLayout()
        self.add_btn = QtGui.QPushButton('Add', self)
        self.add_btn.clicked.connect(self.add_ppt_file)
        self.add_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_file_add.table)))
        self.add_btn.setIconSize(QtCore.QSize(16,16))
        self.open_btn = QtGui.QPushButton('Open', self)
        self.open_btn.clicked.connect(self.open_ppt_file)
        self.open_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.open_btn.setIconSize(QtCore.QSize(16,16))
        open_layout.addWidget(self.open_btn)
        open_layout.addWidget(self.add_btn)
    
        self.ppt_list_table = QtGui.QTableWidget()
        self.ppt_list_table.setColumnCount(4)
        self.ppt_list_table.setHorizontalHeaderItem(0, QtGui.QTableWidgetItem("Name"))
        self.ppt_list_table.setHorizontalHeaderItem(1, QtGui.QTableWidgetItem("Slide"))
        self.ppt_list_table.setHorizontalHeaderItem(2, QtGui.QTableWidgetItem("Path"))
        self.ppt_list_table.setHorizontalHeaderItem(3, QtGui.QTableWidgetItem("Parg"))
    
        header = self.ppt_list_table.horizontalHeader()
        header.setResizeMode(1, QtGui.QHeaderView.ResizeToContents)
        header.setResizeMode(3, QtGui.QHeaderView.ResizeToContents)
        
        self.move_up_btn = QtGui.QPushButton('', self)
        self.move_up_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_arrow_up.table)))
        self.move_up_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.move_up_btn, QtCore.SIGNAL('clicked()'), self.move_item_up)
        
        self.move_down_btn = QtGui.QPushButton('', self)
        self.move_down_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_arrow_down.table)))
        self.move_down_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.move_down_btn, QtCore.SIGNAL('clicked()'), self.move_itme_down)
        
        self.delete_btn = QtGui.QPushButton('', self)
        self.delete_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_delete.table)))
        self.delete_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.delete_btn, QtCore.SIGNAL('clicked()'), self.delete_item)
        
        self.delete_all_btn = QtGui.QPushButton('', self)
        self.delete_all_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_trash.table)))
        self.delete_all_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.delete_all_btn, QtCore.SIGNAL('clicked()'), self.delete_all_item)

        self.sort_asnd_btn = QtGui.QPushButton('', self)
        self.sort_asnd_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_table_sort_asc.table)))
        self.sort_asnd_btn.setIconSize(QtCore.QSize(16,16))
    
        self.sort_dsnd_btn = QtGui.QPushButton('', self)
        self.sort_dsnd_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_slide_setting.table)))
        self.sort_dsnd_btn.setIconSize(QtCore.QSize(16,16))
        self.connect(self.sort_dsnd_btn, QtCore.SIGNAL('clicked()'), self.set_ppt_slide_data)
        
        #self.play = QtGui.QPushButton('', self)
        #self.play.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_play.table)))
        #self.play.setIconSize(QtCore.QSize(16,16))
        
        btn_layout = QtGui.QHBoxLayout()
        btn_layout.addWidget(self.move_up_btn)
        btn_layout.addWidget(self.move_down_btn)
        btn_layout.addWidget(self.delete_btn)
        btn_layout.addWidget(self.delete_all_btn)
        btn_layout.addWidget(self.sort_asnd_btn)
        btn_layout.addWidget(self.sort_dsnd_btn)
        #btn_layout.addWidget(self.play)
        
        publish_layout = QtGui.QGridLayout()
        publish_layout.addWidget(QtGui.QLabel('Date'), 1, 0) 
        self.publish_date  = QtGui.QLineEdit(self)
        publish_layout.addWidget(self.publish_date, 1, 1)
        
        self.add_publish_date = QtGui.QCheckBox()
        self.add_publish_date.stateChanged.connect(self.add_publish_date_state_changed)
        publish_layout.addWidget(self.add_publish_date, 1, 2)
        self.publish_date.setText(datetime.datetime.now().strftime("%Y-%m-%d"))
        self.add_publish_date.setChecked(True)
        
        publish_layout.addWidget(QtGui.QLabel('Name'), 2, 0) 
        self.publish_title = QtGui.QComboBox(self)
        self.publish_title.addItems(const._worship_type)
        #self.publish_title.currentIndexChanged.connect(self.custom_worship_type)
        publish_layout.addWidget(self.publish_title, 2, 1)

        self.user_input_btn = QtGui.QPushButton('', self)
        self.user_input_btn.clicked.connect(self.set_custom_worship_type)
        self.user_input_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_keyboard.table)))
        self.user_input_btn.setIconSize(QtCore.QSize(16,16))
        self.user_input_btn.setToolTip('User worship type')
        publish_layout.addWidget(self.user_input_btn, 2, 2)

        
        publish_layout.addWidget(QtGui.QLabel('Sorc'), 3, 0)
        self.src_directory_path  = QtGui.QLineEdit()
        publish_layout.addWidget(self.src_directory_path, 3, 1)		

        self.copy_src_path_button = QtGui.QPushButton('', self)
        self.copy_src_path_button.clicked.connect(self.copy_srcpath_to_dest)
        self.copy_src_path_button.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_copy_src_path.table)))
        self.copy_src_path_button.setIconSize(QtCore.QSize(16,16))
        self.copy_src_path_button.setToolTip('Copy src path to dest')
        publish_layout.addWidget(self.copy_src_path_button, 3, 2)
    
        publish_layout.addWidget(QtGui.QLabel('Dest'), 4, 0)
        self.save_directory_path  = QtGui.QLineEdit(os.getcwd())
        self.save_directory_button = QtGui.QPushButton('', self)
        self.save_directory_button.clicked.connect(self.get_save_directory_path)
        self.save_directory_button.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_folder_open.table)))
        self.save_directory_button.setIconSize(QtCore.QSize(16,16))
        self.save_directory_button.setToolTip('save folder')
        
        publish_layout.addWidget(self.save_directory_path, 4, 1)
        publish_layout.addWidget(self.save_directory_button, 4, 2)
        
        run_layout = QtGui.QGridLayout()
        isz = 32
        self.run_convert = QtGui.QPushButton('', self)
        self.run_convert.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_convert.table)))
        self.run_convert.setIconSize(QtCore.QSize(isz,isz))
        self.run_convert.clicked.connect(self.create_liveppt)
        self.run_convert.setToolTip('Create Subtitle PPTX')
    
        self.outline_btn = QtGui.QPushButton('', self)
        self.outline_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_outline.table)))
        self.outline_btn.setIconSize(QtCore.QSize(isz,isz))
        self.outline_btn.clicked.connect(self.create_outline_text)
        self.outline_btn.setToolTip('Outline Effect on Subtitle')
    
        self.shadow_btn = QtGui.QPushButton('', self)
        self.shadow_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_shadow.table)))
        self.shadow_btn.setIconSize(QtCore.QSize(isz,isz))
        self.shadow_btn.clicked.connect(self.create_shadow_text)
        self.shadow_btn.setToolTip('Shadow Effect on Subtitle')
                
        self.merge_btn = QtGui.QPushButton('', self)
        #self.merge_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_merge.table)))
        self.merge_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_ppt_txt.table)))
        self.merge_btn.setIconSize(QtCore.QSize(isz,isz))
        #self.merge_btn.clicked.connect(self.run_merge_ppt)
        self.merge_btn.clicked.connect(self.get_txt_from_ppt)
        self.merge_btn.setToolTip('N/A')
        
        self.ppt_pptx_btn = QtGui.QPushButton('', self)
        self.ppt_pptx_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_ppt_pptx.table)))
        self.ppt_pptx_btn.setIconSize(QtCore.QSize(isz,isz))
        self.ppt_pptx_btn.clicked.connect(self.convert_ppt_to_pptx)
        self.ppt_pptx_btn.setToolTip('Convert PPT to PPTX')
    
        self.ppt_img_btn = QtGui.QPushButton('', self)
        self.ppt_img_btn.setIcon(QtGui.QIcon(QtGui.QPixmap(icon_ppt_image.table)))
        self.ppt_img_btn.setIconSize(QtCore.QSize(isz,isz))
        self.ppt_img_btn.clicked.connect(self.convert_ppt_to_image)
        self.ppt_img_btn.setToolTip('Save as Images')
        
        run_layout.addWidget(self.run_convert, 0, 0)
        run_layout.addWidget(self.outline_btn, 0, 1)
        run_layout.addWidget(self.shadow_btn, 0, 2)
        run_layout.addWidget(self.merge_btn, 1, 0)
        run_layout.addWidget(self.ppt_pptx_btn, 1, 1)
        run_layout.addWidget(self.ppt_img_btn, 1, 2)
        
        layout.addRow(open_layout)
        layout.addRow(self.ppt_list_table)
        layout.addRow(btn_layout)
        layout.addRow(publish_layout)
        layout.addRow(run_layout)
        self.ppt_tab.setLayout(layout)
        self.global_message.appendPlainText('... PPT Tab UI created')

    def set_ppt_slide_data(self):
        slide_data_dlg.get_slide_data(self.ppt_subtitle_slide)
        
    def get_ppt_list(self):
        nppt = self.ppt_list_table.rowCount()
        return [os.path.splitext(self.ppt_list_table.item(i,0).text())[0] 
                for i in range(nppt)]
        
    def set_custom_worship_type(self):
        if self.ppt_list_table.rowCount() is 0: return
        
        res = QUserWorshipType(self)
        if res.exec_() is not 1:
            return
            
        src = res.get_source()
        
        if src == 0: # current file on the table
            fn = res.ppt_list.currentText()    
        else: # fx source
            fn = res.user_input.text()
        
        nwt = len(_worship_type)-1
        self.publish_title.setEditable(True)
        self.publish_title.setItemText(nwt, fn)
        self.publish_title.setEditable(False)
        self.publish_title.setCurrentIndex(nwt)
        
    def copy_srcpath_to_dest(self):
        self.save_directory_path.setText(self.src_directory_path.text())
        
    def get_txt_from_ppt(self):
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
        self.global_message.appendPlainText('... PPT to TXT')

        for npr in range(nppt):
            path = self.ppt_list_table.item(npr, 2).text()
            file = self.ppt_list_table.item(npr,0).text()
            ppt_file = os.path.join(path, file)
            txt_file = os.path.join(self.save_directory_path.text(), 
                        "%s.txt"%os.path.splitext(file)[0])
            
            try:
                os.remove(txt_file)
            except OSError as e:
                e_str = str(e)
                if func.access_denied(e_str):
                    e_str += "%s is already opened!"%sfn
                    msgcom.message_box(msgcom.message_error, e_str)
                    continue
            
            fp = open(txt_file, "wt")
            self.global_message.appendPlainText('... Open: %s'%txt_file)
            
            try:
                prs = pptx.Presentation(ppt_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            fp.write("%s\n"%shape.text)
                
            except Exception as e:
                self.global_message.appendPlainText(str(e))
                return
            fp.close()
            self.global_message.appendPlainText('... Close')
        self.global_message.appendPlainText('... Success')
        QtGui.QMessageBox.question(QtGui.QWidget(), 'completed!', 'PPTX to TXT', 
        QtGui.QMessageBox.Yes)
        
    def run_merge_ppt(self):
        
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
    
        try:
            self.global_message.appendPlainText('... Merge PPT: open PowerPoint')
            Application = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
    
        save_file = self.get_save_file(pfix='')
        sfn = os.path.join(self.save_directory_path.text(),save_file)
        
        if os.path.isfile(sfn):
            ans = QtGui.QMessageBox.question(self, 'Continue?', 
                    '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
            if ans == QtGui.QMessageBox.No: return
            else:
                try:
                    os.remove(sfn)
                except OSError as e:
                    if func.access_denied(e_str):
                        e_str += "%s is already opened!"%sfn
                        msgcom.message_box(msgcom.message_error, e_str)
                        self.global_message.appendPlainText(e_str)
                        return 
        try:
            self.global_message.appendPlainText('... Open Presentation')
            # create new presentation object
            # presumes if the file exist, it must be deleted
            dest_ppt = Application.Presentations.Add()
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
            
        dest_ppt.Slides.Add(1, 6)
        index = 1
        self.global_message.appendPlainText('... Start InsertFromFile')
        for npr in range(nppt):
            sp = self.ppt_list_table.item(npr, 2).text()
            sf = self.ppt_list_table.item(npr,0).text()
            src = os.path.join(sp, sf)
            src_ppt = Application.Presentations.Open(src)
            dest_ppt.Slides.InsertFromFile(src, index)
            index += len(src_ppt.Slides)+1
            dest_ppt.Slides.Add(index, 6)
            src_ppt.Close()
        self.global_message.appendPlainText('... End')
        
        self.global_message.appendPlainText('... Start change background color')
        for i, sld in enumerate(dest_ppt.Slides):
            sld.Select()
            # ppViewSlide : 1
            Application.ActiveWindow.ViewType = 1
            Application.ActiveWindow.Activate()
            sr = Application.ActiveWindow.Selection.SlideRange
            sr.FollowMasterBackground = False
            #sr.Background.Fill.Solid = True
            c = self.ppt_subtitle_slide.slide.back_col
            sr.Background.Fill.ForeColor.RGB = RGB(c.r, c.g, c.b)
        self.global_message.appendPlainText('... End')
        
        dest_ppt.SaveAs(sfn)
        Application.Quit()
        self.global_message.appendPlainText('... Merge PPt: success\n')
        
    def create_shadow_text(self, sorce):
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
        
        res = QChooseFxSource()
        if res.exec_() == 1:
            src = res.get_source()
        else: return
        
        if src == 0: # current file on the table
            sfn = os.path.join(self.ppt_list_table.item(0,2).text(),
                                self.ppt_list_table.item(0,0).text())
        else: # fx source
            save_file = self.get_save_file()
            sfn = os.path.join(self.save_directory_path.text(),save_file)
    
        if os.path.isfile(sfn):
            ans = QtGui.QMessageBox.question(self, 'Continue?', 
                    '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
            if ans == QtGui.QMessageBox.No: return
            
        try:
            self.global_message.appendPlainText('... Shadow Text: open PowerPoint')
            Application = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
    
        self.ppt_subtitle_slide.textbox.fx.shadow.Style   = self.fx_shadow_style.currentIndex()+1
        self.ppt_subtitle_slide.textbox.fx.shadow.OffsetX = int(self.fx_shadow_tbl.item(1,1).text())
        self.ppt_subtitle_slide.textbox.fx.shadow.OffsetY = int(self.fx_shadow_tbl.item(2,1).text())
        self.ppt_textbox.fx.shadow.Blur    = int(self.fx_shadow_tbl.item(3,1).text())
        self.ppt_textbox.fx.shadow.Transparency = float(self.fx_shadow_tbl.item(4,1).text())
        self.global_message.appendPlainText(str(self.ppt_textbox.fx.shadow))
        
        try:
            Presentation = Application.Presentations.Open(sfn)
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
            
        for i, sld in enumerate(Presentation.Slides):
            for shp in sld.Shapes:
                sdw = shp.TextFrame2.TextRange.Font.Shadow
                sdw.Visible = 1 # msoTrue
                sdw.Style   = self.ppt_textbox.fx.shadow.Style
                sdw.OffsetX = self.ppt_textbox.fx.shadow.OffsetX
                sdw.OffsetY = self.ppt_textbox.fx.shadow.OffsetY
                sdw.Blur    = self.ppt_textbox.fx.shadow.Blur
                sdw.Transparency = self.ppt_textbox.fx.shadow.Transparency
            
        Presentation.Save()
        self.global_message.appendPlainText('... Shadow Text: success\n')
        Application.Quit()

    def create_outline_text(self, sorce):
    
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
        
        res = QChooseFxSource()
        if res.exec_() == 1:
            src = res.get_source()
        else: return
        
        if src == 0: # current file on the table
            sfn = os.path.join(self.ppt_list_table.item(0,2).text(),
                            self.ppt_list_table.item(0,0).text())
        else: # fx source
            save_file = self.get_save_file()
            sfn = os.path.join(self.save_directory_path.text(),save_file)
            
        if os.path.isfile(sfn):
            ans = QtGui.QMessageBox.question(self, 'Continue?', 
                    '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
            if ans == QtGui.QMessageBox.No: return
        
        try:
            self.global_message.appendPlainText('... Outline Text: open PowerPoint')
            Application = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
    
        try:
            Presentation = Application.Presentations.Open(sfn)
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
        
        
        h = self.fx_outline_show.isChecked()
        c = get_rgb(self.fx_outline_tbl.item(1,1).text())
        s = msoLine.index_to_style(self.fx_outline_style.currentIndex()) # style
        d = self.fx_outline_dash_style.currentIndex()
        t = float(self.fx_outline_tbl.item(4,1).text()) # transprancy
        w = float(self.fx_outline_tbl.item(5,1).text()) # weight
        
        self.ppt_textbox.fx.outline.show = h
        self.ppt_textbox.fx.outline.col = c
        self.ppt_textbox.fx.outline.style = s
        self.ppt_textbox.fx.outline.dash = d
        self.ppt_textbox.fx.outline.transprancy = t
        self.ppt_textbox.fx.outline.weight = w
            
        self.global_message.appendPlainText(str(self.ppt_textbox.fx.outline))
        
        for i, sld in enumerate(Presentation.Slides):
            sld.Select()
            # ppViewSlide : 1
            Application.ActiveWindow.ViewType = 1
            Application.ActiveWindow.Activate()
        
            try:
                sld.Shapes[0].Select()
            except Exception: # in case of an empty slide
                continue 
                
            fnt = Application.ActiveWindow.Selection.TextRange2.Font
            # msoCTrue: 1, msoTrue: -1, msoFalse: 0
            fnt.Line.Visible = h if h is True else False 
            fnt.Line.ForeColor.RGB = RGB(c.r, c.g, c.b)
            fnt.Line.Style = s
            fnt.Line.DashStyle = d
            fnt.Line.Transparency = t
            fnt.Line.Weight = w
        
        Presentation.Save()
        Application.Quit()
        self.global_message.appendPlainText("... Outline text: success\n")
            
    def convert_ppt_to_image(self):
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
    
        res = QImageResolution()
        if res.exec_() == 1:
            w, h = res.get_resolution()
        else: return
            
        h = int(float(w * self.ppt_subtitle_slide.slide.hgt) / self.ppt_subtitle_slide.slide.wid)
                
        try:
            self.global_message.appendPlainText('... Convert PPT to Img: open PowerPoint')
            Application = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText("... Fail: %s"%str(e))
            return
    
        for npr in range(nppt):
            path = self.ppt_list_table.item(npr, 2).text()
            sorc = os.path.join(path, self.ppt_list_table.item(npr,0).text())
            ff = os.path.splitext(self.ppt_list_table.item(npr,0).text())
            img_folder = os.path.join(path, ff[0])
            try:
                os.mkdir(img_folder)
            except FileExistsError:
                pass
    
            Presentation = Application.Presentations.Open(sorc)
            for i, sld in enumerate(Presentation.Slides):
                sld.Select()
                sld.Export(os.path.join(img_folder,"%s%03d.jpg"%(ff[0],i)), "JPG", w, h)
        Application.Quit()
        self.global_message.appendPlainText("... Convert PPT to Img: success\n")
        
        QtGui.QMessageBox.question(QtGui.QWidget(), 'completed!', "%s"%img_folder, 
        QtGui.QMessageBox.Yes)
        
    def convert_ppt_to_pptx(self):
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
    
        try:
            self.global_message.appendPlainText('... PPT to PPTX: open PowerPoint')
            Application = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Open Powerpoint Fail: %s'%str(e))
            return
            
        for npr in range(nppt):
            path = self.ppt_list_table.item(npr, 2).text()
            sorc = os.path.join(path, self.ppt_list_table.item(npr,0).text())
            try:
                Presentation = Application.Presentations.Open(sorc)
            except Exception as e:
                res = QtGui.QMessageBox.question(QtGui.QWidget(), '', "Error:%s"%str(e),
                QtGui.QMessageBox.Cancel)
                self.global_message.appendPlainText('... Open Fiel Fail: %s'%str(e))
                Application.Quit()
                return
                
            fname = os.path.splitext(self.ppt_list_table.item(npr,0).text())
            dest = os.path.join(self.save_directory_path.text(), "%s.pptx"%(fname[0]))
            try:
                Presentation.Saveas(dest)
            except Exception as e:
                res = QtGui.QMessageBox.question(QtGui.QWidget(), 'Continue?', "Error:%s"%str(e),\
                QtGui.QMessageBox.Yes|QtGui.QMessageBox.Cancel)
                if res is QtGui.QMessageBox.Cancel:
                    Presentation.Close()
                    Application.Quit()
                    return
            Presentation.Close()
        Application.Quit()
        self.global_message.appendPlainText("PPTX: %s"%dest)
        self.global_message.appendPlainText("... PPT to PPTX: success\n")
        QtGui.QMessageBox.question(QtGui.QWidget(), 'completed!', "%s"%dest, QtGui.QMessageBox.Yes)
        
    def set_common_var(self):
    
        # TO DO
        # Read default settings : slide size
        slide_size_index = const.get_slide_sizeindex_16x9()
        
        #self.ppt_slide = ppt_slide_info(slide_size_index)
        #self.ppt_slide.back_col = color._color_white
        #self.ppt_textbox = ppt_textbox_info()
        
        self.ppt_subtitle_slide = ppt_slide_data(slide_size_index)
        self.ppt_subtitle_slide.textbox.sy = self.ppt_subtitle_slide.slide.hgt - const._default_txt_hgt
        self.ppt_subtitle_slide.textbox.hgt = const._default_txt_hgt
        self.ppt_subtitle_slide.textbox.fill.solid_col = color._default_solid_fill_color
        #self.ppt_subtitle_slide.textbox.vanchor = const._default_textframe_vanchor_index
        #self.ppt_subtitle_slide.textbox.pp_align = const._default_pragraph_align_index
        
        self.ppt_respread_slide = ppt_slide_data(slide_size_index)
        self.ppt_txtppt_slide   = ppt_slide_data(slide_size_index)
        self.ppt_hymal_slide    = ppt_hymal_data(slide_size_index)
        
        self.ppt_shadow = ppt_shadow_info() 
        self.ppt_fx = ppt_fx_info()
        self.application_path = os.getcwd()
        
    def get_default_pptx(self):
        return os.path.join(self.application_path, const._default_pptx_template)
        
    def custom_worship_type(self):
        cid = self.publish_title.currentIndex()
        nwt = len(_worship_type)-1
        
        if cid == nwt:
            txt, ok = QtGui.QInputDialog.getText(self, 'Custom Worship Type', "Enter")
            if ok:
                self.publish_title.setEditable(True)
                self.publish_title.setItemText(cid, txt)
                self.publish_title.setEditable(False)
        
    def add_publish_date_state_changed(self):
        if self.add_publish_date.isChecked():
            self.publish_date.setEnabled(True)
        else:
            self.publish_date.setEnabled(False)
        
    def get_save_directory_path(self):
        startingDir = os.getcwd() 
        path = QtGui.QFileDialog.getExistingDirectory(None, 'Save folder', startingDir, 
        #path = QtGui.QFileDialog.getExistingDirectory(None, 'Save folder', '', 
        QtGui.QFileDialog.ShowDirsOnly)
        if not path: return
        self.save_directory_path.setText(path)
        os.chdir(path)
    
    def open_ppt_file(self):
        title = self.open_btn.text()
        self.ppt_filenames = QtGui.QFileDialog.getOpenFileNames(self, title, 
        directory=self.src_directory_path.text(), 
        filter="PPTX (*.pptx);;PPT (*.ppt);;All files (*.*)")
        nppt = len(self.ppt_filenames)
    
        if nppt: 
            cur_tab = self.tabs.currentIndex()
            if self.tabs.tabText(cur_tab) == const._ppttab_text:
                self.clear_pptlist_table()
                self.ppt_list_table.setRowCount(nppt)
                for k in range(nppt):
                    fpath, fname = os.path.split(self.ppt_filenames[k])
                    try:
                        prs = pptx.Presentation(self.ppt_filenames[k])
                        nslide = len(prs.slides)
                    except Exception:
                        nslide = 0
                        pass
                            
                    self.ppt_list_table.setItem(k, 0, QtGui.QTableWidgetItem(fname))
                    self.ppt_list_table.setItem(k, 1, QtGui.QTableWidgetItem("%d"%nslide))
                    self.ppt_list_table.setItem(k, 2, QtGui.QTableWidgetItem(fpath))
                    self.ppt_list_table.setItem(k, 3, QtGui.QTableWidgetItem("%d"%const._default_txt_nparagraph))
                self.src_directory_path.setText(fpath)
        
    def add_ppt_file(self):
        title = self.add_btn.text()
        files = QtGui.QFileDialog.getOpenFileNames(self, title, 
        directory=self.src_directory_path.text(), 
        filter="PPTX (*.pptx);;PPT (*.ppt);;All files (*.*)")
        
        nppt = len(files)
        if nppt:
            cur_row = self.ppt_list_table.rowCount()
            for k in range(nppt):
                j = k + cur_row
                fpath, fname = os.path.split(files[k])
                try:
                    prs = pptx.Presentation(self.ppt_filenames[k])
                except Exception as e:
                    res = QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', str(e), QtGui.QMessageBox.Yes|QtGui.QMessageBox.Cancel)
                    if res is QtGui.QMessageBox.Cancel:
                        return
                        
                self.ppt_list_table.insertRow(j)
                self.ppt_list_table.setItem(j, 0, QtGui.QTableWidgetItem(fname))
                self.ppt_list_table.setItem(j, 1, QtGui.QTableWidgetItem("%s"%len(prs.slides)))
                self.ppt_list_table.setItem(j, 2, QtGui.QTableWidgetItem(fpath))
                self.ppt_list_table.setItem(j, 3, QtGui.QTableWidgetItem("%d"%const._default_txt_nparagraph))
            self.src_directory_path.setText(fpath)
    
    #http://stackoverflow.com/questions/9166087/move-row-up-and-down-in-pyqt4
    def move_itme_down(self):
        row = self.ppt_list_table.currentRow()
        column = self.ppt_list_table.currentColumn()
        ncolumn = self.ppt_list_table.columnCount()
        new_pos = row+2
    
        if row < self.ppt_list_table.rowCount()-1:
            self.ppt_list_table.insertRow(new_pos)
            for i in range(ncolumn):
                self.ppt_list_table.setItem(new_pos,i,self.ppt_list_table.takeItem(row,i))
                self.ppt_list_table.setCurrentCell(new_pos,column)
            
            self.ppt_list_table.removeRow(row)        

    def move_item_up(self):    
        row = self.ppt_list_table.currentRow()
        column = self.ppt_list_table.currentColumn();
        ncolumn = self.ppt_list_table.columnCount()
        if row > 0:
            self.ppt_list_table.insertRow(row-1)
            for i in range(ncolumn):
                self.ppt_list_table.setItem(row-1,i,self.ppt_list_table.takeItem(row+1,i))
                self.ppt_list_table.setCurrentCell(row-1,column)
            self.ppt_list_table.removeRow(row+1)        
    
    def clear_pptlist_table(self):
        for i in reversed(range(self.ppt_list_table.rowCount())):
            self.ppt_list_table.removeRow(i)
        self.ppt_list_table.setRowCount(0)
    
    def delete_all_item(self):
        self.clear_pptlist_table()
        
    def delete_item(self): 
        row_count = self.ppt_list_table.rowCount()
        if row_count == 1: self.delete_all_item()
        if row_count > 1:
            column = self.ppt_list_table.currentColumn();
            row = self.ppt_list_table.currentRow()
            for i in range(self.ppt_list_table.columnCount()):
                self.ppt_list_table.setItem(row,i,self.ppt_list_table.takeItem(row+1,i))
                self.ppt_list_table.setCurrentCell(row,column)
            self.ppt_list_table.removeRow(row+1)
            self.ppt_list_table.setRowCount(row_count-1)
    
    def add_empty_slide(self, dest_ppt, layout, bk_col):
        dest_slide = dest_ppt.slides.add_slide(layout)
        bk = dest_slide.background
        bk.fill.solid()
        bk.fill.fore_color.rgb = pptx.dml.color.RGBColor(bk_col.r, bk_col.g, bk_col.b)
        return dest_slide
    
    def slide_has_text(self, slide):
        p_list = []
        for shape in slide.shapes:
            if shape.has_text_frame: 
                for paragraph in shape.text_frame.paragraphs:
                    run_list = []
                    for run in paragraph.runs:
                        run_list.append(run.text)
                    line_text = ''.join(run_list)
                    p_list.append(line_text)
            else: return ''
        txt = ''.join(p_list)
        return txt
    
    def get_save_file(self, pfix = const._liveppt_postfix):
        if self.add_publish_date.isChecked():
            save_file = "%s %s%s.pptx"%(
            self.publish_date.text(), 
            self.publish_title.currentText(),pfix)
        else:
            save_file = "%s%s.pptx"%(
            self.publish_title.currentText(), pfix)
        return save_file	
        
    # dp  : dest_ppt
    # bsl : blank_slide_layout
    # bc  : self.ppt_slide.back_col
    
    def insert_empty_slide(self, dp, slide_data, bsl, bc):
        pbx = self.ppt_subtitle_slide.textbox
        dest_slide = self.add_empty_slide(dp, bsl, bc)
        txt_box = self.add_textbox(dest_slide, slide_data, pbx.sx, pbx.sy, pbx.wid,	pbx.hgt)
        txt_f = txt_box.text_frame
        self.set_textbox(txt_f, MSO_AUTO_SIZE.NONE, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.MIDDLE,
        pbx.left_margin, pbx.top_margin, pbx.right_margin, pbx.bottom_margin)
    
    def create_liveppt(self):
        #import copy
        nppt = self.ppt_list_table.rowCount()
        if nppt == 0: return
    
        save_file = self.get_save_file()
        sfn = os.path.join(self.save_directory_path.text(),save_file)
        
        if os.path.isfile(sfn):
            ans = QtGui.QMessageBox.question(self, 'Continue?', 
                    '%s already exist!'%sfn, QtGui.QMessageBox.Yes, QtGui.QMessageBox.No)
            if ans == QtGui.QMessageBox.No: return
            else:
                try:
                    os.remove(sfn)
                except OSError as e:
                    e_str = str(e)
                    if func.access_denied(e_str):
                        e_str += "%s is already opened!"%sfn
                        msgcom.message_box(msgcom.message_error, e_str)
                        return
            
        self.global_message.appendPlainText('... Create Subtitle PPT')
        try:
            # 2/14/21 py2exe empty presentaion doesn't work
            #dest_ppt = pptx.Presentation(_default_pptx_template)
            dest_ppt = pptx.Presentation(self.get_default_pptx())
        except Exception as e:
            e_str = "Error: can't open pptx document\n%s"%str(e)
            self.global_message.appendPlainText(e_str)
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', e_str)
            return
            
        dest_ppt.slide_width = pptx.util.Inches(self.ppt_subtitle_slide.slide.wid)
        dest_ppt.slide_height = pptx.util.Inches(self.ppt_subtitle_slide.slide.hgt)
        blank_slide_layout = dest_ppt.slide_layouts[6]
    
        for npr in range(nppt):
            sp = self.ppt_list_table.item(npr, 2).text()
            sf = self.ppt_list_table.item(npr,0).text()
            npf = os.path.join(sp, sf)
            self.global_message.appendPlainText("%02d: %s"%(npr+1, sf))
            npg = int(self.ppt_list_table.item(npr, 3).text())
            try:
                src = pptx.Presentation(npf)
            except Exception as e:
                QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "{}".format(e), QtGui.QMessageBox.Yes)
                self.global_message.appendPlainText('... Fail: %s\n'%e)
                dest_ppt = None
                return
    
            pbx = self.ppt_subtitle_slide.textbox
            
            for slide in src.slides:
                in_text = self.slide_has_text(slide)
                if in_text=='' and self.ppt_subtitle_slide.slide.deep_copy:
                    self.insert_empty_slide(dest_ppt, self.ppt_subtitle_slide, blank_slide_layout, self.ppt_subtitle_slide.slide.back_col)
                    continue
                    
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    p_list = []
                    for paragraph in shape.text_frame.paragraphs:
                        run_list = []
                        for run in paragraph.runs:
                            run_list.append(run.text)
                        line_text = ''.join(run_list)
                        
                        # Delete 1. 2. 3. ...
                        match = func._find_lyric_number.search(line_text)
                        if match:
                            line_text = func._find_lyric_number.sub('', line_text)
                        
                        # skip hymal chapter info. ex: (찬송기 123장)
                        line_text = line_text.strip()
                        match = func._skip_hymal_info.search(line_text)
                        if match or not line_text: 
                            continue
                        p_list.append(line_text)
                    
                    np_list = len(p_list)
    
                    for j in range(0, np_list, npg):
                        if npg > 1 and (np_list-j) < npg: break
                        dest_slide = self.add_empty_slide(dest_ppt, blank_slide_layout, self.ppt_subtitle_slide.slide.back_col)
                        txt_box = self.add_textbox(dest_slide, self.ppt_subtitle_slide, pbx.sx, pbx.sy, pbx.wid, pbx.hgt)
                        txt_f = txt_box.text_frame
                        self.set_textbox(txt_f, 
                                        MSO_AUTO_SIZE.NONE, 
                                        const.get_textframe_vanchortype(self.ppt_subtitle_slide.textbox.vanchor),
                                        const.get_paragraph_aligntype(self.ppt_subtitle_slide.textbox.pp_align),         #MSO_ANCHOR.MIDDLE, 
                                        #MSO_ANCHOR.MIDDLE,
                                        pbx.left_margin, pbx.top_margin,
                                        pbx.right_margin, pbx.bottom_margin)           
            
                        k_list = p_list[j:j+npg]
                        if pbx.word_wrap:
                            self.add_paragraph(txt_f, k_list[0], True)
                            for w in k_list[1:]:
                                self.add_paragraph(txt_f, w)
                        else:
                            self.add_paragraph(txt_f, ' '.join(k_list), True)
    
                    if npg == 1: continue
                    left_over = np_list%npg
                    if left_over:
                        dest_slide = self.add_empty_slide(dest_ppt, blank_slide_layout, self.ppt_subtitle_slide.slide.back_col)
                        txt_box = self.add_textbox(dest_slide, self.ppt_subtitle_slide, pbx.sx, pbx.sy, pbx.wid, pbx.hgt)
                        txt_f = txt_box.text_frame
                        self.set_textbox(txt_f,MSO_AUTO_SIZE.NONE, 
                                            #MSO_ANCHOR.MIDDLE, 
                                            #MSO_ANCHOR.MIDDLE,
                                            const.get_textframe_vanchortype(self.ppt_subtitle_slide.textbox.vanchor),
                                            const.get_paragraph_aligntype(self.ppt_subtitle_slide.textbox.pp_align),
                                            pbx.left_margin, pbx.top_margin,
                                            pbx.right_margin, pbx.bottom_margin)
                                              
                        k_list = p_list[j:j+left_over]
                        if pbx.word_wrap:
                            self.add_paragraph(txt_f, k_list[0], True)
                            for w in k_list[1:]:
                                self.add_paragraph(txt_f, w)
                        else:
                            self.add_paragraph(txt_f, ' '.join(k_list), True)
            self.insert_empty_slide(dest_ppt, self.ppt_subtitle_slide, blank_slide_layout, self.ppt_subtitle_slide.slide.back_col)
            
        try:
            dest_ppt.save(sfn)
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "{}".format(e), QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
        
        #if self.ppt_subtitle_slide.textbox.fill.show: #textbox_fill.isChecked():
        #    #ft = self.textbox_fill_type.currentIndex()
        #    ft = self.ppt_subtitle_slide.textbox.fill.fill_type
        #    if ft is const._TEXTBOX_GRADIENT_FILL:
        #        self.fill_textbox(sfn, bool(ft))
            #else:
            #    self.fill_textbox(sfn, bool(ft))
    
        QtGui.QMessageBox.question(QtGui.QWidget(), 'Completed!', sfn, QtGui.QMessageBox.Yes)
        self.global_message.appendPlainText('Dest: %s\n%s\n... Create Subtitle PPT: success\n'%(
        save_file, str(pbx)))
    
    def add_textbox(self, ds, slide_data, sx, sy, wid, hgt):
        txt_box = ds.shapes.add_textbox(\
                pptx.util.Inches(sx),
                pptx.util.Inches(sy),
                pptx.util.Inches(wid),
                pptx.util.Inches(hgt))
        
        #if bkc:
        if slide_data.textbox.fill.show:
            if slide_data.textbox.fill.type == const._TEXTBOX_SOLID_FILL:
                sc = slide_data.textbox.fill.solid_col
                txt_box.fill.solid()
                txt_box.fill.fore_color.rgb = pptx.dml.color.RGBColor(sc.r, sc.g, sc.b)
            else:
                txt_box.fill.gradient()
                gradient_stops = txt_box.fill.gradient_stops
                st1 = gradient_stops[0]
                st2 = gradient_stops[1]
                gc1 = slide_data.textbox.fill.gradient_col1
                gc2 = slide_data.textbox.fill.gradient_col2
                st1.color.rgb = pptx.dml.color.RGBColor(gc1.r, gc1.g, gc1.b)
                st2.color.rgb = pptx.dml.color.RGBColor(gc2.r, gc2.g, gc2.b)

        return txt_box
        
    def add_paragraph(self, tf, txt, first=False):
        if first:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = txt
        self.set_paragraph(p, func.get_textalign(self.ppt_subtitle_slide.textbox.pp_align),
                            self.ppt_subtitle_slide.textbox.font_name,
                            self.ppt_subtitle_slide.textbox.font_size,
                            self.ppt_subtitle_slide.textbox.font_col,
                            self.ppt_subtitle_slide.textbox.font_bold)
        
    # slide(az,va,ha): MSO_AUTO_SIZE.NONE, MSO_ANCHOR.BOTTOM, MSO_ANCHOR.MIDDLE
    # hymal(az,va,ha): MSO_AUTO_SIZE.NONE, MSO_ANCHOR.BOTTOM, MSO_ANCHOR.LEFT
    def set_textbox(self, tf, az, va, ha, lm, tm, rm, bm, wr=True):
        tf.auto_size = az
        tf.vertical_anchor = va
        tf.horizontal_anchor= ha
        tf.margin_left   = pptx.util.Inches(lm)
        tf.margin_top    = pptx.util.Inches(tm)
        tf.margin_right  = pptx.util.Inches(rm)
        tf.margin_bottom = pptx.util.Inches(bm)
        tf.word_wrap = wr
    
    # slide(al,fn,fz,fc,bl) = PP_ALIGN.CENTER, 
    #                         self.ppt_subtitle_slide.textbox.font_name,
    #						  self.ppt_subtitle_slide.textbox.font_size,
    #                         self.ppt_subtitle_slide.textbox.font_col,
    #                         self.ppt_subtitle_slide.textbox.font_bold
    # hymal(al,fn,fz,fc,bl) = PP_ALIGN.LEFT, 
    def set_paragraph(self, p, al, fn, fz, fc, bl=True):
        p.alignment = al
        p.font.name = fn
        p.font.size = pptx.util.Pt(fz)
        p.font.color.rgb = pptx.dml.color.RGBColor(fc.r, fc.g, fc.b)
        p.font.bold = bl
    
    '''
    Public Sub gradient_fill()
    Dim sld As Slide
    Dim shp As Shape
    
    For Each sld In ActivePresentation.Slides
        For Each shp In sld.Shapes
            shp.Fill.Visible = msoCTrue
            shp.Fill.GradientAngle = 0
            shp.Fill.GradientStops.Insert RGB(0, 0, 0), 0, 0
            shp.Fill.GradientStops.Insert RGB(64, 64, 64), 0.46, 0
            shp.Fill.GradientStops.Insert RGB(217, 217, 217), 0.77, 0.7
            shp.Fill.GradientStops.Insert RGB(255, 255, 255), 1, 1
        Next shp
    Next sld
    End Sub
    '''
    # https://docs.microsoft.com/en-us/office/vba/api/powerpoint.fillformat.twocolorgradient
    # https://docs.microsoft.com/en-us/office/vba/api/office.msogradientstyle
    
    def fill_textbox(self, sfn, gradient=True):
    
        try:
            self.global_message.appendPlainText('... Fill textbox: open PowerPoint')
            Application = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
    
        try:
            Presentation = Application.Presentations.Open(sfn)
        except Exception as e:
            QtGui.QMessageBox.question(QtGui.QWidget(), 'Error', "%s"%str(e), 
            QtGui.QMessageBox.Yes)
            self.global_message.appendPlainText('... Fail: %s'%str(e))
            return
            
        self.global_message.appendPlainText('Type: %s'%('Gradient' if gradient else 'Solid'))
        sc = get_rgb(self.textbox_solid_color.text())
        g1 = get_rgb(self.textbox_gradient_color1.text())
        g2 = get_rgb(self.textbox_gradient_color2.text())
        
        self.global_message.appendPlainText('S.Col : %s\nG.Col1: %s\nG.Col2: %s'%(str(sc), str(g1), str(g2)))
        for i, sld in enumerate(Presentation.Slides):
            for shp in sld.Shapes:
                shp.Fill.Visible = True
                if gradient:
                    shp.Fill.TwoColorGradient(2, 1)
                    shp.Fill.GradientStops[0].Color.RGB = RGB(g1.r, g1.g ,g1.b )
                    shp.Fill.GradientStops[1].Color.RGB = RGB(g2.r, g2.g ,g2.b )
                else:
                    shp.Fill.ForeColor.RGB = RGB(sc.r,sc.g,sc.b)
                
                # work on Powerpoint VBS but not win32com
                #shp.Fill.GradientAngle = 0
                #shp.Fill.GradientStops.Insert(RGB(0  , 0  , 0  ), 0   , 0  )
                #shp.Fill.GradientStops.Insert(RGB(64 , 64 , 64 ), 0.46, 0  )
                #shp.Fill.GradientStops.Insert(RGB(217, 217, 217), 0.77, 0.7)
                #shp.Fill.GradientStops.Insert(RGB(255, 255, 255), 1   , 1  )
            
        Presentation.Save()
        Application.Quit()
        self.global_message.appendPlainText('... Gradient Fill: success')

def main():
    app = QtGui.QApplication(sys.argv)
    #QtGui.QApplication.setStyle(QtGui.QStyleFactory.create(u'Motif'))
    #QtGui.QApplication.setStyle(QtGui.QStyleFactory.create(u'CDE'))
    QtGui.QApplication.setStyle(QtGui.QStyleFactory.create(u'Plastique'))
    #QtGui.QApplication.setStyle(QtGui.QStyleFactory.create(u'Cleanlooks'))
    lppt = QLivePPT()
    sys.exit(app.exec_())

if __name__ == '__main__':
    main()	
